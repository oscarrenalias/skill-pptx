"""Verification checks for .pptx files."""

from __future__ import annotations

import math
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# ── Namespaces ────────────────────────────────────────────────────────────────

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ── Check 1: unfilled placeholder hint text ───────────────────────────────────

_PLACEHOLDER_PATTERNS: tuple[str, ...] = (
    "click to add title",
    "click to add text",
    "click to add subtitle",
    "click to edit master title style",
    "click to edit master text styles",
    "add title",
    "add text",
    "enter text here",
    "place subtitle here",
)


def _check_unfilled_placeholders(
    slide_index: int,
    slide,
    errors: list[str],
) -> None:
    """Check 1: detect shapes whose text matches an unfilled placeholder pattern."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if not text.strip():
            continue
        if text.strip().lower() in _PLACEHOLDER_PATTERNS:
            errors.append(
                f"Slide {slide_index}: '{shape.name}' has unfilled placeholder"
                f' text \u2014 "{text.strip()}"'
            )


# ── Check 2: font size below minimum ─────────────────────────────────────────

_MIN_FONT_UNITS = 1200  # 12pt in hundredths-of-a-point


def _check_font_sizes(
    slide_index: int,
    slide,
    slide_height: int,
    warnings: list[str],
) -> None:
    """Check 2: append a warning for each font below 12pt; skip footer-region shapes."""
    footer_threshold = int(slide_height * 0.9)

    for shape in slide.shapes:
        # Skip shapes in the footer region (top > 90% of slide height)
        top = shape.top
        if top is not None and top > footer_threshold:
            continue

        # Walk a:rPr[@sz] elements in this shape's XML
        sp_elem = shape._element
        for rpr in sp_elem.iter(f"{{{NS_A}}}rPr"):
            sz = rpr.get("sz")
            if sz is None:
                continue
            try:
                sz_int = int(sz)
            except ValueError:
                continue
            if sz_int < _MIN_FONT_UNITS:
                # Convert from hundredths-of-a-point to pt for the message
                pt = sz_int / 100.0
                warnings.append(
                    f"Slide {slide_index}: '{shape.name}' font {pt:.1f}pt"
                    f" is below minimum (12pt)"
                )


# ── Check 3: shape overflow ───────────────────────────────────────────────────

_OVERFLOW_TOLERANCE = 50_000  # EMU — ignore sub-pixel rounding
_EMU_PER_INCH = 914_400


def _check_shape_overflow(
    slide_index: int,
    slide,
    slide_width: int,
    slide_height: int,
    errors: list[str],
) -> None:
    """Check 3: detect shapes whose bounding box overflows slide dimensions."""
    footer_threshold = int(slide_height * 0.9)

    for shape in slide.shapes:
        # Skip table shapes
        if shape.has_table:
            continue

        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height

        if left is None or top is None or width is None or height is None:
            continue

        # Skip footer-region shapes (top > 90% of slide height)
        if top > footer_threshold:
            continue

        left = int(left)
        top = int(top)
        width = int(width)
        height = int(height)

        # Negative position → overflows that edge
        if left < 0:
            errors.append(
                f"Slide {slide_index}: '{shape.name}' has negative left"
                f" position ({left} EMU) — shape overflows left edge"
            )
        if top < 0:
            errors.append(
                f"Slide {slide_index}: '{shape.name}' has negative top"
                f" position ({top} EMU) — shape overflows top edge"
            )

        # Right overflow: allow up to _OVERFLOW_TOLERANCE past slide boundary
        if left + width > slide_width + _OVERFLOW_TOLERANCE:
            errors.append(
                f"Slide {slide_index}: '{shape.name}' overflows right edge"
                f" (left+width={left + width}, slide_width={slide_width})"
            )

        # Bottom overflow
        if top + height > slide_height + _OVERFLOW_TOLERANCE:
            errors.append(
                f"Slide {slide_index}: '{shape.name}' overflows bottom edge"
                f" (top+height={top + height}, slide_height={slide_height})"
            )


# ── Check 4: text clipping ────────────────────────────────────────────────────

_DEFAULT_FONT_EMU = 18 * 12700  # 18pt in EMU (1pt = 12700 EMU)
_LINE_SPACING = 1.35
_AVG_CHAR_WIDTH_FACTOR = 0.55


def _para_font_size_emu(para) -> int:
    """Return effective font size for a paragraph in EMU; falls back to 18pt."""
    for run in para.runs:
        if run.font.size is not None:
            return int(run.font.size)
    return _DEFAULT_FONT_EMU


def _check_text_clipping(
    slide_index: int,
    slide,
    slide_height: int,
    errors: list[str],
    warnings: list[str],
) -> None:
    """Check 4: estimate rendered text height and warn/error on likely clipping."""
    footer_threshold = int(slide_height * 0.9)

    for shape in slide.shapes:
        # Skip table shapes and non-text shapes
        if shape.has_table or not shape.has_text_frame:
            continue

        top = shape.top
        if top is None:
            continue

        # Skip footer-region shapes (top > 90% of slide height)
        if int(top) > footer_threshold:
            continue

        tf = shape.text_frame

        # Skip shapes with no text
        if not tf.text.strip():
            continue

        box_width = shape.width
        box_height = shape.height
        if box_width is None or box_height is None or box_width <= 0 or box_height <= 0:
            continue

        box_width = int(box_width)
        box_height = int(box_height)
        word_wrap = tf.word_wrap  # True, False, or None (None treated as True)

        estimated_height = 0.0
        for para in tf.paragraphs:
            font_size_emu = _para_font_size_emu(para)

            if word_wrap is False:
                line_count = 1
            else:
                avg_char_width = _AVG_CHAR_WIDTH_FACTOR * font_size_emu
                chars_per_line = box_width / avg_char_width if avg_char_width > 0 else 1
                text_len = len(para.text)
                line_count = max(1, math.ceil(text_len / chars_per_line)) if chars_per_line > 0 else 1

            estimated_height += font_size_emu * _LINE_SPACING * line_count

        if estimated_height <= 0:
            continue

        overage = (estimated_height - box_height) / box_height

        if overage > 0.20:
            est_in = estimated_height / _EMU_PER_INCH
            box_in = box_height / _EMU_PER_INCH
            pct = round(overage * 100)
            errors.append(
                f"Slide {slide_index}: '{shape.name}' text may be clipped"
                f" (est={est_in:.2f}\" vs box={box_in:.2f}\", +{pct}%)"
            )
        elif overage > 0.08:
            est_in = estimated_height / _EMU_PER_INCH
            box_in = box_height / _EMU_PER_INCH
            pct = round(overage * 100)
            warnings.append(
                f"Slide {slide_index}: '{shape.name}' text may be clipped"
                f" (est={est_in:.2f}\" vs box={box_in:.2f}\", +{pct}%)"
            )


# ── Check 5: significant shape overlap ───────────────────────────────────────

_TINY_SHAPE_THRESHOLD = 50_000  # EMU — skip shapes smaller than this
_OVERLAP_AREA_THRESHOLD = 0.3 * _EMU_PER_INCH * _EMU_PER_INCH  # 0.3 sq in in EMU²
_OVERLAP_FRACTION_THRESHOLD = 0.20  # 20% of the smaller shape's area


def _shape_area(shape) -> int:
    """Return the area of a shape in EMU²."""
    w = shape.width
    h = shape.height
    if w is None or h is None:
        return 0
    return int(w) * int(h)


def _intersection_area(a, b) -> int:
    """Return the intersection area (in EMU²) of two shapes' bounding boxes."""
    a_left = int(a.left or 0)
    a_top = int(a.top or 0)
    a_right = a_left + int(a.width or 0)
    a_bottom = a_top + int(a.height or 0)

    b_left = int(b.left or 0)
    b_top = int(b.top or 0)
    b_right = b_left + int(b.width or 0)
    b_bottom = b_top + int(b.height or 0)

    inter_left = max(a_left, b_left)
    inter_top = max(a_top, b_top)
    inter_right = min(a_right, b_right)
    inter_bottom = min(a_bottom, b_bottom)

    if inter_right <= inter_left or inter_bottom <= inter_top:
        return 0
    return (inter_right - inter_left) * (inter_bottom - inter_top)


def _fully_contains(outer, inner) -> bool:
    """Return True if outer's bounding box fully contains inner's."""
    o_left = int(outer.left or 0)
    o_top = int(outer.top or 0)
    o_right = o_left + int(outer.width or 0)
    o_bottom = o_top + int(outer.height or 0)

    i_left = int(inner.left or 0)
    i_top = int(inner.top or 0)
    i_right = i_left + int(inner.width or 0)
    i_bottom = i_top + int(inner.height or 0)

    return o_left <= i_left and o_top <= i_top and o_right >= i_right and o_bottom >= i_bottom


def _check_shape_overlap(
    slide_index: int,
    slide,
    warnings: list[str],
) -> None:
    """Check 5: warn when two non-tiny shapes overlap significantly."""
    shapes = list(slide.shapes)
    for i in range(len(shapes)):
        a = shapes[i]
        # Skip tiny shapes
        w_a = a.width
        h_a = a.height
        if w_a is None or h_a is None:
            continue
        if int(w_a) < _TINY_SHAPE_THRESHOLD or int(h_a) < _TINY_SHAPE_THRESHOLD:
            continue

        for j in range(i + 1, len(shapes)):
            b = shapes[j]
            w_b = b.width
            h_b = b.height
            if w_b is None or h_b is None:
                continue
            if int(w_b) < _TINY_SHAPE_THRESHOLD or int(h_b) < _TINY_SHAPE_THRESHOLD:
                continue

            inter = _intersection_area(a, b)
            if inter <= 0:
                continue

            # Skip if one fully contains the other
            if _fully_contains(a, b) or _fully_contains(b, a):
                continue

            area_a = _shape_area(a)
            area_b = _shape_area(b)
            smaller_area = min(area_a, area_b)
            if smaller_area <= 0:
                continue

            # Both thresholds must be exceeded
            if inter > _OVERLAP_AREA_THRESHOLD and inter > _OVERLAP_FRACTION_THRESHOLD * smaller_area:
                inter_sq_in = inter / (_EMU_PER_INCH * _EMU_PER_INCH)
                pct = round(100 * inter / smaller_area)
                warnings.append(
                    f"Slide {slide_index}: '{a.name}' and '{b.name}' overlap"
                    f" significantly ({inter_sq_in:.2f} sq in, {pct}% of smaller shape)"
                )


# ── Public API ────────────────────────────────────────────────────────────────


def verify_pptx(path: Path) -> dict:
    """Run quality checks on a .pptx file.

    Returns a dict with keys:
        errors      -- list of error message strings
        warnings    -- list of warning message strings
        slide_count -- total number of slides
    """
    prs = Presentation(Path(path))
    slide_width: int = int(prs.slide_width or Emu(9144000))   # default 10in
    slide_height: int = int(prs.slide_height or Emu(6858000))  # default 7.5in

    errors: list[str] = []
    warnings: list[str] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        _check_unfilled_placeholders(slide_index, slide, errors)
        _check_font_sizes(slide_index, slide, slide_height, warnings)
        _check_shape_overflow(slide_index, slide, slide_width, slide_height, errors)
        _check_text_clipping(slide_index, slide, slide_height, errors, warnings)
        _check_shape_overlap(slide_index, slide, warnings)

    return {
        "errors": errors,
        "warnings": warnings,
        "slide_count": len(prs.slides),
    }
