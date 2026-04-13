from pathlib import Path

from pptx import Presentation


def extract_text(path: Path, slides: list[int] | None = None) -> str:
    """Extract text from a .pptx file.

    Args:
        path: Path to the .pptx file.
        slides: Optional 1-based list of slide indices to include. If None,
                all slides are included.

    Returns:
        A string containing the extracted text, with ``--- Slide N ---``
        delimiter lines separating each slide's content.
    """
    prs = Presentation(path)
    parts: list[str] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        if slides is not None and slide_index not in slides:
            continue

        parts.append(f"--- Slide {slide_index} ---")

        sorted_shapes = sorted(
            slide.shapes,
            key=lambda s: (s.top if s.top is not None else 0,
                           s.left if s.left is not None else 0),
        )

        for shape in sorted_shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                line = "".join(run.text for run in paragraph.runs)
                if line:
                    parts.append(line)

    return "\n".join(parts)
