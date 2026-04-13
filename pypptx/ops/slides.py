"""Read and write operations for slides."""

from __future__ import annotations

import io
import re
import shutil
import tempfile
import xml.etree.ElementTree as StdET
import zipfile
from contextlib import contextmanager
from pathlib import Path
from typing import Generator

import defusedxml.ElementTree as ET
from pptx import Presentation

from pypptx.ops.clean import clean_unused_files
from pypptx.ops.pack import pack, unpack

# ── Namespaces ────────────────────────────────────────────────────────────────

NS_PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_RELS = "http://schemas.openxmlformats.org/package/2006/relationships"
NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

_REL_TYPE_SLIDE = f"{NS_R}/slide"
_REL_TYPE_SLIDE_LAYOUT = f"{NS_R}/slideLayout"
_REL_TYPE_NOTES_SLIDE = f"{NS_R}/notesSlide"
_CT_SLIDE = (
    "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
)


# ── Presentation helper ───────────────────────────────────────────────────────


def _open_presentation(path: Path) -> Presentation:
    """Open a Presentation from a .pptx file or an unpacked directory."""
    if path.is_dir():
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            # OPC spec: [Content_Types].xml must be the first ZIP entry.
            ct = path / "[Content_Types].xml"
            if ct.exists():
                zf.write(ct, "[Content_Types].xml")
            for f in sorted(path.rglob("*")):
                if not f.is_file():
                    continue
                arcname = str(f.relative_to(path))
                if arcname == "[Content_Types].xml":
                    continue  # already written first
                zf.write(f, arcname)
        buf.seek(0)
        return Presentation(buf)
    return Presentation(path)


# ── XML helpers ───────────────────────────────────────────────────────────────


def _register_ns_from_file(path: Path) -> None:
    """Register all XML namespace declarations found in *path* with StdET."""
    content = path.read_bytes()
    for prefix, uri in re.findall(rb'xmlns:(\w+)\s*=\s*"([^"]+)"', content):
        StdET.register_namespace(prefix.decode(), uri.decode())
    for uri in re.findall(rb'xmlns="([^"]+)"', content):
        StdET.register_namespace("", uri.decode())


def _get_ordered_slides(pptx_dir: Path) -> list[dict]:
    """Return ordered list of slide info from sldIdLst.

    Each dict contains:
        rid      -- rId string in presentation.xml.rels
        sld_id   -- id attribute value (string) from sldIdLst
        target   -- target as written in presentation.xml.rels
        filename -- bare filename, e.g. 'slide1.xml'
    """
    prs_xml = pptx_dir / "ppt" / "presentation.xml"
    prs_rels = pptx_dir / "ppt" / "_rels" / "presentation.xml.rels"

    rid_to_target: dict[str, str] = {}
    if prs_rels.exists():
        tree = ET.parse(str(prs_rels))
        for rel in tree.getroot().findall(f"{{{NS_RELS}}}Relationship"):
            if rel.get("Type", "") == _REL_TYPE_SLIDE:
                rid_to_target[rel.get("Id", "")] = rel.get("Target", "")

    result = []
    if prs_xml.exists():
        tree = ET.parse(str(prs_xml))
        sld_id_lst = tree.getroot().find(f"{{{NS_PML}}}sldIdLst")
        if sld_id_lst is not None:
            for el in sld_id_lst:
                rid = el.get(f"{{{NS_R}}}id", "")
                sld_id = el.get("id", "")
                target = rid_to_target.get(rid, "")
                filename = target.rsplit("/", 1)[-1] if target else ""
                result.append(
                    {"rid": rid, "sld_id": sld_id, "target": target, "filename": filename}
                )
    return result


def _get_sorted_layouts(pptx_dir: Path) -> list[str]:
    """Return sorted list of layout filenames from ppt/slideLayouts/."""
    layouts_dir = pptx_dir / "ppt" / "slideLayouts"
    if not layouts_dir.exists():
        return []
    return sorted(
        f.name
        for f in layouts_dir.iterdir()
        if f.is_file() and f.suffix == ".xml"
    )


def _next_slide_filename(pptx_dir: Path) -> str:
    """Return the next available slideN.xml filename."""
    slides_dir = pptx_dir / "ppt" / "slides"
    existing = {
        f.name
        for f in slides_dir.iterdir()
        if f.is_file() and f.name.startswith("slide") and f.name.endswith(".xml")
    }
    n = 1
    while f"slide{n}.xml" in existing:
        n += 1
    return f"slide{n}.xml"


def _next_presentation_rid(prs_rels_path: Path) -> str:
    """Return the next available rId in presentation.xml.rels."""
    tree = ET.parse(str(prs_rels_path))
    rids: set[int] = set()
    for rel in tree.getroot().findall(f"{{{NS_RELS}}}Relationship"):
        rid = rel.get("Id", "")
        if rid.startswith("rId"):
            try:
                rids.add(int(rid[3:]))
            except ValueError:
                pass
    n = 1
    while n in rids:
        n += 1
    return f"rId{n}"


def _next_sld_id_value(prs_xml_path: Path) -> int:
    """Return the next available sldId id value (min 256)."""
    tree = ET.parse(str(prs_xml_path))
    root = tree.getroot()
    sld_id_lst = root.find(f"{{{NS_PML}}}sldIdLst")
    max_id = 255
    if sld_id_lst is not None:
        for el in sld_id_lst:
            try:
                val = int(el.get("id", "0"))
                if val > max_id:
                    max_id = val
            except ValueError:
                pass
    return max_id + 1


def _add_content_type_override(pptx_dir: Path, part_name: str) -> None:
    """Add an Override entry for *part_name* in [Content_Types].xml."""
    ct_path = pptx_dir / "[Content_Types].xml"
    _register_ns_from_file(ct_path)
    tree = ET.parse(str(ct_path))
    root = tree.getroot()
    if not part_name.startswith("/"):
        part_name = "/" + part_name
    override = StdET.SubElement(root, f"{{{NS_CT}}}Override")
    override.set("PartName", part_name)
    override.set("ContentType", _CT_SLIDE)
    StdET.register_namespace("", NS_CT)
    StdET.ElementTree(root).write(str(ct_path), xml_declaration=True, encoding="UTF-8")


def _add_presentation_rel(pptx_dir: Path, new_rid: str, target: str) -> None:
    """Add a slide Relationship to ppt/_rels/presentation.xml.rels."""
    prs_rels = pptx_dir / "ppt" / "_rels" / "presentation.xml.rels"
    _register_ns_from_file(prs_rels)
    tree = ET.parse(str(prs_rels))
    root = tree.getroot()
    rel = StdET.SubElement(root, f"{{{NS_RELS}}}Relationship")
    rel.set("Id", new_rid)
    rel.set("Type", _REL_TYPE_SLIDE)
    rel.set("Target", target)
    StdET.register_namespace("", NS_RELS)
    StdET.ElementTree(root).write(str(prs_rels), xml_declaration=True, encoding="UTF-8")


def _append_sld_id_entry(pptx_dir: Path, new_id: int, new_rid: str) -> None:
    """Append a new <p:sldId> entry to sldIdLst in presentation.xml."""
    prs_xml = pptx_dir / "ppt" / "presentation.xml"
    _register_ns_from_file(prs_xml)
    tree = ET.parse(str(prs_xml))
    root = tree.getroot()
    sld_id_lst = root.find(f"{{{NS_PML}}}sldIdLst")
    if sld_id_lst is None:
        raise RuntimeError("sldIdLst not found in presentation.xml")
    el = StdET.SubElement(sld_id_lst, f"{{{NS_PML}}}sldId")
    el.set("id", str(new_id))
    el.set(f"{{{NS_R}}}id", new_rid)
    StdET.register_namespace("p", NS_PML)
    StdET.register_namespace("r", NS_R)
    StdET.ElementTree(root).write(str(prs_xml), xml_declaration=True, encoding="UTF-8")


def _strip_notes_rels(rels_path: Path) -> None:
    """Remove notesSlide relationships from *rels_path* in-place."""
    if not rels_path.exists():
        return
    _register_ns_from_file(rels_path)
    tree = ET.parse(str(rels_path))
    root = tree.getroot()
    to_remove = [
        rel
        for rel in root.findall(f"{{{NS_RELS}}}Relationship")
        if rel.get("Type", "") == _REL_TYPE_NOTES_SLIDE
    ]
    for rel in to_remove:
        root.remove(rel)
    StdET.register_namespace("", NS_RELS)
    StdET.ElementTree(root).write(str(rels_path), xml_declaration=True, encoding="UTF-8")


def _blank_slide_xml() -> str:
    """Return minimal blank slide XML (layout reference goes in .rels)."""
    return (
        "<?xml version='1.0' encoding='UTF-8' standalone='yes'?>\n"
        f'<p:sld xmlns:a="{NS_A}" xmlns:p="{NS_PML}" xmlns:r="{NS_R}">'
        "<p:cSld>"
        "<p:spTree>"
        "<p:nvGrpSpPr>"
        '<p:cNvPr id="1" name=""/>'
        "<p:cNvGrpSpPr/>"
        "<p:nvPr/>"
        "</p:nvGrpSpPr>"
        "<p:grpSpPr>"
        "<a:xfrm>"
        '<a:off x="0" y="0"/>'
        '<a:ext cx="0" cy="0"/>'
        '<a:chOff x="0" y="0"/>'
        '<a:chExt cx="0" cy="0"/>'
        "</a:xfrm>"
        "</p:grpSpPr>"
        "</p:spTree>"
        "</p:cSld>"
        "<p:clrMapOvr>"
        "<a:masterClr/>"
        "</p:clrMapOvr>"
        "</p:sld>"
    )


def _slide_rels_xml(layout_target: str) -> str:
    """Return a .rels XML for a new slide referencing layout at *layout_target*."""
    return (
        "<?xml version='1.0' encoding='UTF-8' standalone='yes'?>\n"
        f'<Relationships xmlns="{NS_RELS}">'
        f'<Relationship Id="rId1" Type="{_REL_TYPE_SLIDE_LAYOUT}"'
        f' Target="{layout_target}"/>'
        "</Relationships>"
    )


# ── Core write helpers ────────────────────────────────────────────────────────


def _add_slide_to_dir(
    pptx_dir: Path,
    *,
    duplicate: int | None = None,
    layout: int | None = None,
) -> dict:
    """Add a slide to an unpacked PPTX directory; returns {file, index}."""
    prs_xml = pptx_dir / "ppt" / "presentation.xml"
    prs_rels = pptx_dir / "ppt" / "_rels" / "presentation.xml.rels"
    slides_dir = pptx_dir / "ppt" / "slides"
    rels_dir = slides_dir / "_rels"

    ordered = _get_ordered_slides(pptx_dir)
    new_filename = _next_slide_filename(pptx_dir)
    new_slide_path = slides_dir / new_filename
    new_rels_path = rels_dir / (new_filename + ".rels")
    new_rid = _next_presentation_rid(prs_rels)
    new_id = _next_sld_id_value(prs_xml)

    if duplicate is not None:
        if duplicate < 1 or duplicate > len(ordered):
            raise ValueError(
                f"Slide index {duplicate} out of range (1-{len(ordered)})"
            )
        src = ordered[duplicate - 1]
        src_slide_path = slides_dir / src["filename"]
        src_rels_path = rels_dir / (src["filename"] + ".rels")

        shutil.copy2(src_slide_path, new_slide_path)
        rels_dir.mkdir(parents=True, exist_ok=True)
        if src_rels_path.exists():
            shutil.copy2(src_rels_path, new_rels_path)
            _strip_notes_rels(new_rels_path)

    else:  # layout mode
        layout_filenames = _get_sorted_layouts(pptx_dir)
        if layout < 1 or layout > len(layout_filenames):  # type: ignore[operator]
            raise ValueError(
                f"Layout index {layout} out of range (1-{len(layout_filenames)})"
            )
        layout_filename = layout_filenames[layout - 1]  # type: ignore[index]
        # Target is relative to the owning part (ppt/slides/slideN.xml)
        layout_rel_target = f"../slideLayouts/{layout_filename}"

        new_slide_path.write_text(_blank_slide_xml(), encoding="UTF-8")
        rels_dir.mkdir(parents=True, exist_ok=True)
        new_rels_path.write_text(_slide_rels_xml(layout_rel_target), encoding="UTF-8")

    # Register in presentation.xml.rels (target relative to ppt/presentation.xml)
    _add_presentation_rel(pptx_dir, new_rid, f"slides/{new_filename}")
    # Register in [Content_Types].xml
    _add_content_type_override(pptx_dir, f"/ppt/slides/{new_filename}")
    # Append to sldIdLst
    _append_sld_id_entry(pptx_dir, new_id, new_rid)

    return {"file": new_filename, "index": len(ordered) + 1}


def _delete_slide_from_dir(pptx_dir: Path, index: int) -> dict:
    """Remove the slide at 1-based *index* from an unpacked directory."""
    ordered = _get_ordered_slides(pptx_dir)
    if index < 1 or index > len(ordered):
        raise ValueError(f"Slide index {index} out of range (1-{len(ordered)})")

    slide_info = ordered[index - 1]
    rid = slide_info["rid"]
    filename = slide_info["filename"]

    slides_dir = pptx_dir / "ppt" / "slides"
    slide_path = slides_dir / filename
    rels_path = slides_dir / "_rels" / (filename + ".rels")
    prs_xml = pptx_dir / "ppt" / "presentation.xml"
    prs_rels = pptx_dir / "ppt" / "_rels" / "presentation.xml.rels"
    ct_path = pptx_dir / "[Content_Types].xml"

    # Remove from sldIdLst in presentation.xml
    _register_ns_from_file(prs_xml)
    prs_tree = ET.parse(str(prs_xml))
    prs_root = prs_tree.getroot()
    sld_id_lst = prs_root.find(f"{{{NS_PML}}}sldIdLst")
    if sld_id_lst is not None:
        for el in [e for e in sld_id_lst if e.get(f"{{{NS_R}}}id") == rid]:
            sld_id_lst.remove(el)
    StdET.register_namespace("p", NS_PML)
    StdET.register_namespace("r", NS_R)
    StdET.ElementTree(prs_root).write(str(prs_xml), xml_declaration=True, encoding="UTF-8")

    # Remove from presentation.xml.rels
    _register_ns_from_file(prs_rels)
    rels_tree = ET.parse(str(prs_rels))
    rels_root = rels_tree.getroot()
    for el in [
        r for r in rels_root.findall(f"{{{NS_RELS}}}Relationship")
        if r.get("Id") == rid
    ]:
        rels_root.remove(el)
    StdET.register_namespace("", NS_RELS)
    StdET.ElementTree(rels_root).write(str(prs_rels), xml_declaration=True, encoding="UTF-8")

    # Remove from [Content_Types].xml
    _register_ns_from_file(ct_path)
    ct_tree = ET.parse(str(ct_path))
    ct_root = ct_tree.getroot()
    part_name = f"/ppt/slides/{filename}"
    for el in [
        e for e in ct_root.findall(f"{{{NS_CT}}}Override")
        if e.get("PartName") == part_name
    ]:
        ct_root.remove(el)
    StdET.register_namespace("", NS_CT)
    StdET.ElementTree(ct_root).write(str(ct_path), xml_declaration=True, encoding="UTF-8")

    # Delete the slide file and its .rels
    slide_path.unlink(missing_ok=True)
    rels_path.unlink(missing_ok=True)

    return {"deleted_file": filename, "deleted_index": index}


def _move_slide_in_dir(pptx_dir: Path, from_idx: int, to_idx: int) -> dict:
    """Reorder sldIdLst moving slide at *from_idx* to *to_idx* (1-based)."""
    ordered = _get_ordered_slides(pptx_dir)
    n = len(ordered)
    if from_idx < 1 or from_idx > n:
        raise ValueError(f"from_idx {from_idx} out of range (1-{n})")
    if to_idx < 1 or to_idx > n:
        raise ValueError(f"to_idx {to_idx} out of range (1-{n})")

    filename = ordered[from_idx - 1]["filename"]

    if from_idx == to_idx:
        return {"file": filename, "from": from_idx, "to": to_idx}

    prs_xml = pptx_dir / "ppt" / "presentation.xml"
    _register_ns_from_file(prs_xml)
    tree = ET.parse(str(prs_xml))
    root = tree.getroot()
    sld_id_lst = root.find(f"{{{NS_PML}}}sldIdLst")
    if sld_id_lst is None:
        raise RuntimeError("sldIdLst not found in presentation.xml")

    elements = list(sld_id_lst)
    for el in elements:
        sld_id_lst.remove(el)

    el = elements.pop(from_idx - 1)
    elements.insert(to_idx - 1, el)

    for el in elements:
        sld_id_lst.append(el)

    StdET.register_namespace("p", NS_PML)
    StdET.register_namespace("r", NS_R)
    StdET.ElementTree(root).write(str(prs_xml), xml_declaration=True, encoding="UTF-8")

    return {"file": filename, "from": from_idx, "to": to_idx}


# ── Context manager ───────────────────────────────────────────────────────────


@contextmanager
def pptx_edit(path: Path) -> Generator[Path, None, None]:
    """Context manager for editing a .pptx file in-place.

    Unpacks *path* to a temporary directory on entry and yields that directory.
    On clean exit, cleans unused files then repacks atomically over the original.
    On exception, the original file is left untouched.

    Parameters
    ----------
    path:
        Path to a .pptx file.
    """
    path = Path(path)
    tmp_dir = Path(tempfile.mkdtemp(prefix="pypptx_edit_"))
    try:
        unpack(path, tmp_dir)
        yield tmp_dir
        clean_unused_files(tmp_dir)
        pack(tmp_dir, path)
    except Exception:
        raise
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ── Public API ────────────────────────────────────────────────────────────────


def list_slides(path: Path) -> list[dict]:
    """Return slides in presentation order.

    Each dict contains:
        index  -- 1-based position in the presentation
        file   -- bare filename, e.g. "slide1.xml"
        hidden -- True only when the slide's show attribute is explicitly 'false' or '0'
    """
    prs = _open_presentation(Path(path))
    result = []
    for i, slide in enumerate(prs.slides, start=1):
        filename = slide.part.partname.rsplit("/", 1)[-1]
        show_attr = slide._element.get("show")
        hidden = show_attr in ("0", "false") if show_attr is not None else False
        result.append({"index": i, "file": filename, "hidden": hidden})
    return result


def list_layouts(path: Path) -> list[dict]:
    """Return all slide layouts from ppt/slideLayouts/ in filename order.

    Each dict contains:
        index -- 1-based, assigned after sorting by filename
        file  -- bare filename, e.g. "slideLayout1.xml"
    """
    prs = _open_presentation(Path(path))
    seen: set[str] = set()
    filenames: list[str] = []
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            filename = layout.part.partname.rsplit("/", 1)[-1]
            if filename not in seen:
                seen.add(filename)
                filenames.append(filename)
    filenames.sort()
    return [{"index": i, "file": f} for i, f in enumerate(filenames, start=1)]


def add_slide(
    path: Path,
    *,
    duplicate: int | None = None,
    layout: int | None = None,
) -> dict:
    """Add a slide to a .pptx file or unpacked directory.

    Exactly one of *duplicate* or *layout* must be supplied.

    Parameters
    ----------
    path:
        Path to a .pptx file or unpacked directory.
    duplicate:
        1-based index of the slide to duplicate. The notes-slide relationship
        is stripped from the copy's .rels file.
    layout:
        1-based index of the layout (as returned by ``list_layouts``) to use
        for a new blank slide.

    Returns
    -------
    Dict with ``file`` (bare filename) and ``index`` (1-based position).

    Raises
    ------
    ValueError
        If both or neither of *duplicate*/*layout* are provided, or if the
        index is out of range.
    """
    if (duplicate is None) == (layout is None):
        raise ValueError("Exactly one of 'duplicate' or 'layout' must be provided.")

    path = Path(path)
    if path.is_file() and path.suffix.lower() == ".pptx":
        with pptx_edit(path) as tmp_dir:
            return _add_slide_to_dir(tmp_dir, duplicate=duplicate, layout=layout)
    return _add_slide_to_dir(path, duplicate=duplicate, layout=layout)


def delete_slide(path: Path, index: int) -> dict:
    """Remove the slide at 1-based *index*.

    Parameters
    ----------
    path:
        Path to a .pptx file or unpacked directory.
    index:
        1-based index of the slide to remove.

    Returns
    -------
    Dict with ``deleted_file`` and ``deleted_index``.

    Raises
    ------
    ValueError
        If *index* is out of range.
    """
    path = Path(path)
    if path.is_file() and path.suffix.lower() == ".pptx":
        with pptx_edit(path) as tmp_dir:
            return _delete_slide_from_dir(tmp_dir, index)
    return _delete_slide_from_dir(path, index)


def move_slide(path: Path, from_idx: int, to_idx: int) -> dict:
    """Reorder slides by moving the slide at *from_idx* to *to_idx* (1-based).

    Parameters
    ----------
    path:
        Path to a .pptx file or unpacked directory.
    from_idx:
        1-based index of the slide to move.
    to_idx:
        1-based destination index.

    Returns
    -------
    Dict with ``file`` (bare filename), ``from``, and ``to``.

    Raises
    ------
    ValueError
        If either index is out of range.
    """
    path = Path(path)
    if path.is_file() and path.suffix.lower() == ".pptx":
        with pptx_edit(path) as tmp_dir:
            return _move_slide_in_dir(tmp_dir, from_idx, to_idx)
    return _move_slide_in_dir(path, from_idx, to_idx)
