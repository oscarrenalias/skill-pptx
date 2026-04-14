import json
import sys
from pathlib import Path
from typing import Callable

import click

import pypptx


def output_result(data: dict, plain: bool, plain_fn: Callable[[dict], str]) -> None:
    """Write data to stdout as JSON or as plain text via plain_fn."""
    if plain:
        sys.stdout.write(plain_fn(data) + "\n")
    else:
        sys.stdout.write(json.dumps(data) + "\n")


@click.group()
@click.version_option(version=pypptx.__version__, prog_name="pypptx")
def cli() -> None:
    """pypptx — PowerPoint manipulation toolkit."""


@cli.command("extract-text")
@click.argument("file", type=click.Path(exists=True, path_type=Path))
@click.option("--slides", "slides_str", default=None, help="Comma-separated 1-based slide indices.")
@click.option("--output", "output_file", default=None, type=click.Path(), help="Write extracted text to this file.")
@click.option("--plain", is_flag=True, default=False, help="Output plain text instead of JSON (applies only with --output).")
def extract_text_cmd(file: Path, slides_str: str | None, output_file: str | None, plain: bool) -> None:
    """Extract text from a .pptx file."""
    from pptx import Presentation

    from pypptx.ops.extract import extract_text

    slides: list[int] | None = None
    if slides_str is not None:
        slides = []
        for token in slides_str.split(","):
            token = token.strip()
            try:
                slides.append(int(token))
            except ValueError:
                click.echo(f"Error: invalid slide index {token!r}", err=True)
                sys.exit(1)

        prs = Presentation(file)
        total = len(prs.slides)
        for idx in slides:
            if idx < 1 or idx > total:
                click.echo(f"Error: slide index {idx} out of range (1-{total})", err=True)
                sys.exit(1)

    try:
        text = extract_text(file, slides=slides)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)

    if output_file is None:
        sys.stdout.write(text + "\n")
    else:
        try:
            Path(output_file).write_text(text, encoding="utf-8")
        except Exception as e:
            click.echo(f"Error writing to {output_file!r}: {e}", err=True)
            sys.exit(1)

        slide_count = sum(1 for line in text.splitlines() if line.startswith("--- Slide "))
        output_result(
            {"output_file": output_file, "slide_count": slide_count},
            plain,
            lambda d: d["output_file"],
        )


@cli.command("unpack")
@click.argument("file", type=click.Path(exists=True, path_type=Path))
@click.argument("output_dir", required=False, default=None, type=click.Path(path_type=Path))
@click.option("--plain", is_flag=True, default=False, help="Output plain text instead of JSON.")
def unpack_cmd(file: Path, output_dir: Path | None, plain: bool) -> None:
    """Unpack a .pptx file into a directory."""
    from pypptx.ops.pack import unpack

    if output_dir is None:
        output_dir = Path(file.stem)

    try:
        result_dir = unpack(file, output_dir)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)

    output_result(
        {"unpacked_dir": str(result_dir)},
        plain,
        lambda d: d["unpacked_dir"],
    )


@cli.command("pack")
@click.argument("src_dir", type=click.Path(exists=True, path_type=Path))
@click.argument("output_file", type=click.Path(path_type=Path))
@click.option("--plain", is_flag=True, default=False, help="Output plain text instead of JSON.")
def pack_cmd(src_dir: Path, output_file: Path, plain: bool) -> None:
    """Pack an unpacked directory back into a .pptx file."""
    from pypptx.ops.pack import pack

    try:
        result_file = pack(src_dir, output_file)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)

    output_result(
        {"output_file": str(result_file)},
        plain,
        lambda d: d["output_file"],
    )


@cli.command("clean")
@click.argument("path", type=click.Path(exists=True, path_type=Path))
@click.option("--plain", is_flag=True, default=False, help="Output plain text instead of JSON.")
def clean_cmd(path: Path, plain: bool) -> None:
    """Remove orphaned files from a .pptx file or unpacked directory."""
    from pypptx.ops.clean import clean_unused_files

    try:
        removed = clean_unused_files(path)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)

    output_result(
        {"removed": removed},
        plain,
        lambda d: "\n".join(d["removed"]),
    )


@cli.command("thumbnails")
@click.argument("file", type=click.Path(exists=True, path_type=Path))
@click.option("--output", "output_prefix", default="thumbnails", show_default=True, help="Output filename prefix.")
@click.option("--cols", default=3, show_default=True, type=int, help="Number of grid columns (max: 6).")
@click.option("--plain", is_flag=True, default=False, help="Output one file path per line instead of JSON.")
def thumbnails_cmd(file: Path, output_prefix: str, cols: int, plain: bool) -> None:
    """Generate labeled thumbnail grid images from a .pptx file."""
    import tempfile

    from pypptx.ops.thumbnails import check_dependencies, generate_thumbnails

    if cols > 6:
        click.echo("Error: --cols must be 6 or fewer.", err=True)
        sys.exit(1)

    check_dependencies()

    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            output_paths = generate_thumbnails(file, output_prefix, temp_dir, cols=cols)
        except Exception as e:
            click.echo(f"Error: {e}", err=True)
            sys.exit(1)

    output_result(
        {"files": [str(p) for p in output_paths]},
        plain,
        lambda d: "\n".join(d["files"]),
    )


@cli.group()
def slide() -> None:
    """Commands for working with slides."""


@slide.command("list")
@click.argument("path", type=click.Path(exists=True, path_type=Path))
@click.option("--plain", is_flag=True, default=False, help="Output plain text instead of JSON.")
def slide_list_cmd(path: Path, plain: bool) -> None:
    """List slides in a .pptx file or unpacked directory."""
    from pypptx.ops.slides import list_slides

    try:
        slides = list_slides(path)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)

    output_result(
        {"slides": slides},
        plain,
        lambda d: "\n".join(
            s["file"] + (" [hidden]" if s["hidden"] else "") for s in d["slides"]
        ),
    )


@slide.command("layouts")
@click.argument("path", type=click.Path(exists=True, path_type=Path))
@click.option("--plain", is_flag=True, default=False, help="Output plain text instead of JSON.")
def slide_layouts_cmd(path: Path, plain: bool) -> None:
    """List slide layouts in a .pptx file or unpacked directory."""
    from pypptx.ops.slides import list_layouts

    try:
        layouts = list_layouts(path)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)

    output_result(
        {"layouts": layouts},
        plain,
        lambda d: "\n".join(l["file"] for l in d["layouts"]),
    )


@slide.command("add")
@click.argument("file", type=click.Path(exists=True, path_type=Path))
@click.option("--duplicate", default=None, type=int, help="1-based index of slide to duplicate.")
@click.option("--layout", default=None, type=int, help="1-based index of layout for a new blank slide.")
@click.option("--position", default=None, type=int, help="1-based insertion position (default: end of presentation).")
@click.option("--plain", is_flag=True, default=False, help="Output plain text instead of JSON.")
def slide_add_cmd(file: Path, duplicate: int | None, layout: int | None, position: int | None, plain: bool) -> None:
    """Add a slide to a .pptx file or unpacked directory."""
    from pypptx.ops.slides import add_slide, move_slide

    if (duplicate is None) == (layout is None):
        click.echo("Error: exactly one of --duplicate or --layout must be provided.", err=True)
        sys.exit(1)

    try:
        result = add_slide(file, duplicate=duplicate, layout=layout)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)

    added_filename = result["file"]
    final_position = result["index"]

    if position is not None and position != final_position:
        try:
            move_result = move_slide(file, final_position, position)
            final_position = move_result["to"]
        except Exception as e:
            click.echo(f"Error: {e}", err=True)
            sys.exit(1)

    output_result(
        {"added_file": added_filename, "position": final_position},
        plain,
        lambda d: f"{d['added_file']} at position {d['position']}",
    )


@slide.command("delete")
@click.argument("file", type=click.Path(exists=True, path_type=Path))
@click.argument("index", type=int)
@click.option("--plain", is_flag=True, default=False, help="Output plain text instead of JSON.")
def slide_delete_cmd(file: Path, index: int, plain: bool) -> None:
    """Delete a slide from a .pptx file or unpacked directory."""
    from pypptx.ops.slides import delete_slide

    try:
        result = delete_slide(file, index)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)

    output_result(
        result,
        plain,
        lambda d: f"deleted {d['deleted_file']} (index {d['deleted_index']})",
    )


@slide.command("move")
@click.argument("file", type=click.Path(exists=True, path_type=Path))
@click.argument("from_idx", type=int, metavar="FROM")
@click.argument("to_idx", type=int, metavar="TO")
@click.option("--plain", is_flag=True, default=False, help="Output plain text instead of JSON.")
def slide_move_cmd(file: Path, from_idx: int, to_idx: int, plain: bool) -> None:
    """Move a slide within a .pptx file or unpacked directory."""
    from pypptx.ops.slides import move_slide

    try:
        result = move_slide(file, from_idx, to_idx)
    except Exception as e:
        click.echo(f"Error: {e}", err=True)
        sys.exit(1)

    output_result(
        result,
        plain,
        lambda d: f"{d['file']}: {d['from']} -> {d['to']}",
    )
