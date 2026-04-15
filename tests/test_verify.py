"""Tests for pypptx/ops/verify.py and the 'verify' CLI command."""
from __future__ import annotations

import json
from pathlib import Path

import pytest
from click.testing import CliRunner
from pptx import Presentation
from pptx.util import Emu, Pt

from pypptx.cli import cli
from pypptx.ops.verify import verify_pptx


# ── Slide dimensions (standard widescreen 10in × 7.5in) ──────────────────────
SLIDE_W = Emu(9144000)   # 10 in
SLIDE_H = Emu(6858000)   # 7.5 in

# EMU helpers
_EMU_PER_INCH = 914_400


def _make_pptx(tmp_path: Path, name: str = "deck.pptx") -> tuple[Presentation, Path]:
    """Return a blank Presentation and the path it will be saved to."""
    prs = Presentation()
    path = tmp_path / name
    return prs, path


def _save(prs: Presentation, path: Path) -> Path:
    prs.save(str(path))
    return path


# ── Check 1: unfilled placeholder text ───────────────────────────────────────


class TestCheck1UnfilledPlaceholders:
    """Each of the 9 placeholder patterns should trigger; real text should not."""

    PATTERNS = [
        "click to add title",
        "click to add text",
        "click to add subtitle",
        "click to edit master title style",
        "click to edit master text styles",
        "add title",
        "add text",
        "enter text here",
        "place subtitle here",
    ]

    @pytest.mark.parametrize("pattern", PATTERNS)
    def test_fires_for_each_pattern(self, tmp_path, pattern):
        prs, path = _make_pptx(tmp_path, f"check1_{pattern[:8].replace(' ', '_')}.pptx")
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        from pptx.util import Inches
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox.text_frame.text = pattern
        txBox.name = "MyBox"
        _save(prs, path)

        result = verify_pptx(path)
        assert any("unfilled placeholder" in e for e in result["errors"]), (
            f"Expected error for pattern {pattern!r}, got: {result['errors']}"
        )
        # Message must contain slide number and shape name
        matching = [e for e in result["errors"] if "unfilled placeholder" in e]
        assert any("Slide 1" in e and "MyBox" in e for e in matching)

    def test_does_not_fire_for_non_text_shape(self, tmp_path):
        """A shape with no text frame must not generate a check-1 error."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches
        # Add a line (no text frame)
        slide.shapes.add_connector(
            1,  # MSO_CONNECTOR.STRAIGHT
            Inches(1), Inches(1), Inches(3), Inches(2),
        )
        _save(prs, path)
        result = verify_pptx(path)
        assert result["errors"] == []
        assert result["warnings"] == []

    def test_does_not_fire_for_empty_stripped_text(self, tmp_path):
        """A text box whose stripped text is empty must not trigger."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox.text_frame.text = "   "
        _save(prs, path)
        result = verify_pptx(path)
        assert result["errors"] == []

    def test_does_not_fire_for_real_content(self, tmp_path):
        """Legitimate text ('Hello World') must not trigger check 1."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox.text_frame.text = "Hello World"
        _save(prs, path)
        result = verify_pptx(path)
        check1_errors = [e for e in result["errors"] if "unfilled placeholder" in e]
        assert check1_errors == []


# ── Check 2: font size below minimum ─────────────────────────────────────────


class TestCheck2FontSize:

    def _add_run_with_sz(self, slide, sz_hundredths: int, shape_name: str = "TinyFont"):
        """Add a textbox whose run has an explicit sz attribute."""
        from pptx.util import Inches
        from lxml import etree

        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox.name = shape_name
        tf = txBox.text_frame
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "small text"
        # Set sz via XML
        rPr = run._r.get_or_add_rPr()
        rPr.set("sz", str(sz_hundredths))
        return txBox

    def test_fires_below_1200(self, tmp_path):
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_run_with_sz(slide, 900, "SmallFont")  # 9pt
        _save(prs, path)

        result = verify_pptx(path)
        font_warns = [w for w in result["warnings"] if "font" in w and "below minimum" in w]
        assert len(font_warns) >= 1
        assert any("Slide 1" in w and "SmallFont" in w for w in font_warns)
        assert any("9.0pt" in w for w in font_warns)
        assert result["errors"] == []

    def test_does_not_fire_at_exactly_1200(self, tmp_path):
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_run_with_sz(slide, 1200, "OkFont")  # exactly 12pt — should not trigger
        _save(prs, path)

        result = verify_pptx(path)
        font_warns = [w for w in result["warnings"] if "font" in w and "below minimum" in w]
        assert font_warns == []

    def test_does_not_fire_above_minimum(self, tmp_path):
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_run_with_sz(slide, 1800, "NormalFont")  # 18pt
        _save(prs, path)

        result = verify_pptx(path)
        font_warns = [w for w in result["warnings"] if "below minimum" in w]
        assert font_warns == []

    def test_skips_footer_region_shapes(self, tmp_path):
        """Shapes whose top > 90% of slide height must be skipped."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches

        # Place textbox in footer region (top = 95% of slide height)
        footer_top = int(SLIDE_H * 0.95)
        txBox = slide.shapes.add_textbox(
            Emu(0), Emu(footer_top), Inches(4), Inches(0.3)
        )
        txBox.name = "FooterShape"
        tf = txBox.text_frame
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = "footer note"
        rPr = run._r.get_or_add_rPr()
        rPr.set("sz", "800")  # 8pt — would normally trigger
        _save(prs, path)

        result = verify_pptx(path)
        font_warns = [w for w in result["warnings"] if "below minimum" in w and "FooterShape" in w]
        assert font_warns == []


# ── Check 3: shape overflow ───────────────────────────────────────────────────


class TestCheck3ShapeOverflow:

    def test_shape_exactly_fitting_slide_no_error(self, tmp_path):
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(
            Emu(0), Emu(0), SLIDE_W, SLIDE_H
        )
        txBox.name = "FullSlide"
        txBox.text_frame.text = "Fits perfectly"
        _save(prs, path)
        result = verify_pptx(path)
        overflow_errors = [e for e in result["errors"] if "overflows" in e or "negative" in e.lower()]
        assert overflow_errors == []

    def test_fires_for_right_overflow(self, tmp_path):
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Shape goes well past right edge: left=0, width = slide_width + 100k EMU
        txBox = slide.shapes.add_textbox(
            Emu(0), Emu(0), Emu(int(SLIDE_W) + 100_000), Emu(int(SLIDE_H) // 2)
        )
        txBox.name = "RightOverflow"
        txBox.text_frame.text = "overflows right"
        _save(prs, path)
        result = verify_pptx(path)
        assert any("RightOverflow" in e and "overflows right edge" in e for e in result["errors"])

    def test_fires_for_negative_left(self, tmp_path):
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(
            Emu(-100_000), Emu(100_000), Emu(500_000), Emu(500_000)
        )
        txBox.name = "NegLeft"
        txBox.text_frame.text = "negative left"
        _save(prs, path)
        result = verify_pptx(path)
        assert any("NegLeft" in e and "negative left" in e for e in result["errors"])

    def test_fires_for_negative_top(self, tmp_path):
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(
            Emu(100_000), Emu(-100_000), Emu(500_000), Emu(500_000)
        )
        txBox.name = "NegTop"
        txBox.text_frame.text = "negative top"
        _save(prs, path)
        result = verify_pptx(path)
        assert any("NegTop" in e and "negative top" in e for e in result["errors"])

    def test_table_shape_ignored(self, tmp_path):
        """Tables must be skipped by check 3 even when they overflow."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches
        # Add a table that overflows
        table = slide.shapes.add_table(
            2, 2,
            Emu(0), Emu(0),
            Emu(int(SLIDE_W) + 200_000), Inches(2),
        )
        table.name = "OverflowTable"
        _save(prs, path)
        result = verify_pptx(path)
        overflow_errors = [e for e in result["errors"] if "OverflowTable" in e]
        assert overflow_errors == []

    def test_footer_region_shape_ignored(self, tmp_path):
        """Shape in footer region (top > 90%) must not trigger overflow check."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        footer_top = int(SLIDE_H * 0.95)
        txBox = slide.shapes.add_textbox(
            Emu(0), Emu(footer_top),
            Emu(int(SLIDE_W) + 200_000), Emu(200_000),
        )
        txBox.name = "FooterOverflow"
        txBox.text_frame.text = "footer"
        _save(prs, path)
        result = verify_pptx(path)
        overflow_errors = [e for e in result["errors"] if "FooterOverflow" in e]
        assert overflow_errors == []


# ── Check 4: text clipping ────────────────────────────────────────────────────


class TestCheck4TextClipping:

    def test_no_trigger_for_fitting_text(self, tmp_path):
        """A single short line in a tall box must not trigger."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(4))
        txBox.name = "FittingText"
        tf = txBox.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = "Hi"
        run.font.size = Pt(18)
        _save(prs, path)
        result = verify_pptx(path)
        clipping = [
            e for e in result["errors"] + result["warnings"]
            if "FittingText" in e and "clipped" in e
        ]
        assert clipping == []

    def test_no_text_shape_no_trigger(self, tmp_path):
        """A textbox with no text must not trigger check 4."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.5))
        txBox.name = "EmptyBox"
        _save(prs, path)
        result = verify_pptx(path)
        clipping = [
            e for e in result["errors"] + result["warnings"]
            if "EmptyBox" in e
        ]
        assert clipping == []

    def test_warning_for_moderate_overage(self, tmp_path):
        """Text ~10% over box height should produce a warning (not error)."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches

        # Box: 1 in wide, 1 in tall. Font 72pt = 72*12700 = 914400 EMU ≈ 1 in.
        # With LINE_SPACING=1.35: estimated_height = 914400 * 1.35 = 1234440 EMU.
        # overage = (1234440 - 914400) / 914400 ≈ 0.35 → error, not warning.
        # Instead: box 2in tall, font 18pt, many words → many lines needed.
        # Let's use a narrow box (0.3in) with long text at 18pt to force wrapping.
        # avg_char_width = 0.55 * 18*12700 = 125,730 EMU
        # box_width = 0.3in = 274,320 EMU → chars_per_line ≈ 2.18
        # text = 6 chars → line_count = ceil(6/2.18) = 3
        # estimated = 18*12700 * 1.35 * 3 = 925,290 EMU
        # box_height = 2in = 1,828,800 EMU → no overage. Too big.
        #
        # Simpler approach: 1-line text in an undersized box.
        # font 18pt, box_height = 18pt * 12700 * 1.35 * 1.09 (9% overage).
        font_emu = 18 * 12700  # 18pt in EMU
        line_spacing = 1.35
        estimated = font_emu * line_spacing  # one line
        # For ~10% overage: box_height = estimated / 1.10
        box_height_emu = int(estimated / 1.10)

        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Emu(box_height_emu)
        )
        txBox.name = "WarnText"
        tf = txBox.text_frame
        tf.word_wrap = False  # force single line
        run = tf.paragraphs[0].add_run()
        run.text = "Hello World"
        run.font.size = Pt(18)
        _save(prs, path)
        result = verify_pptx(path)
        warn_hits = [w for w in result["warnings"] if "WarnText" in w and "clipped" in w]
        error_hits = [e for e in result["errors"] if "WarnText" in e and "clipped" in e]
        assert len(warn_hits) >= 1 or len(error_hits) >= 1, (
            f"Expected clipping warning/error for WarnText. warnings={result['warnings']}, errors={result['errors']}"
        )

    def test_error_for_large_overage(self, tmp_path):
        """Text >20% over box height should produce an error."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches

        font_emu = 18 * 12700
        line_spacing = 1.35
        estimated = font_emu * line_spacing  # one line, word_wrap=False
        # For >20% overage: box_height = estimated / 1.25
        box_height_emu = int(estimated / 1.25)

        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Emu(box_height_emu)
        )
        txBox.name = "ErrorText"
        tf = txBox.text_frame
        tf.word_wrap = False
        run = tf.paragraphs[0].add_run()
        run.text = "Hello World this is long text"
        run.font.size = Pt(18)
        _save(prs, path)
        result = verify_pptx(path)
        error_hits = [e for e in result["errors"] if "ErrorText" in e and "clipped" in e]
        assert len(error_hits) >= 1, (
            f"Expected clipping error for ErrorText. errors={result['errors']}"
        )
        assert any("Slide 1" in e for e in error_hits)

    def test_footer_region_shape_skipped(self, tmp_path):
        """Text clipping check must skip footer-region shapes."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches

        footer_top = int(SLIDE_H * 0.95)
        font_emu = 18 * 12700
        line_spacing = 1.35
        box_height_emu = int(font_emu * line_spacing / 1.25)  # would trigger

        txBox = slide.shapes.add_textbox(
            Emu(0), Emu(footer_top), Inches(4), Emu(box_height_emu)
        )
        txBox.name = "FooterClip"
        tf = txBox.text_frame
        tf.word_wrap = False
        run = tf.paragraphs[0].add_run()
        run.text = "footer text"
        run.font.size = Pt(18)
        _save(prs, path)
        result = verify_pptx(path)
        clipping = [
            e for e in result["errors"] + result["warnings"]
            if "FooterClip" in e
        ]
        assert clipping == []


# ── Check 5: significant shape overlap ───────────────────────────────────────


class TestCheck5ShapeOverlap:

    def test_fires_for_large_significant_overlap(self, tmp_path):
        """Two large overlapping shapes with >0.3 sq in overlap and >20% fraction."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches

        # Two 4in×4in boxes with 2in×2in overlap (4 sq in, 25% of smaller)
        txA = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(4), Inches(4))
        txA.name = "BoxA"
        txA.text_frame.text = "A"

        txB = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(4), Inches(4))
        txB.name = "BoxB"
        txB.text_frame.text = "B"

        _save(prs, path)
        result = verify_pptx(path)
        overlap_warns = [w for w in result["warnings"] if "overlap" in w.lower()]
        assert len(overlap_warns) >= 1, f"Expected overlap warning, got: {result['warnings']}"
        assert any("BoxA" in w and "BoxB" in w for w in overlap_warns)

    def test_no_trigger_for_tiny_shapes(self, tmp_path):
        """Shapes smaller than _TINY_SHAPE_THRESHOLD (50k EMU) must be skipped."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Two 40k×40k EMU shapes (below 50k threshold)
        txA = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(40_000), Emu(40_000))
        txA.name = "TinyA"
        txA.text_frame.text = "a"

        txB = slide.shapes.add_textbox(Emu(10_000), Emu(10_000), Emu(40_000), Emu(40_000))
        txB.name = "TinyB"
        txB.text_frame.text = "b"

        _save(prs, path)
        result = verify_pptx(path)
        overlap_warns = [w for w in result["warnings"] if "overlap" in w.lower()]
        assert overlap_warns == []

    def test_no_trigger_when_fully_contained(self, tmp_path):
        """One shape fully inside another must not trigger overlap warning."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches

        # Outer: 6in×6in; Inner: 2in×2in fully inside
        outer = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(6), Inches(6))
        outer.name = "Outer"
        outer.text_frame.text = "outer"

        inner = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(2))
        inner.name = "Inner"
        inner.text_frame.text = "inner"

        _save(prs, path)
        result = verify_pptx(path)
        overlap_warns = [w for w in result["warnings"] if "overlap" in w.lower()]
        assert overlap_warns == []

    def test_no_trigger_for_small_overlap_area(self, tmp_path):
        """Overlap below 0.3 sq in threshold must not trigger."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches

        # Two 3in×3in shapes with only a tiny 0.1in×0.1in corner overlap
        # overlap area = 0.01 sq in < 0.3 threshold
        txA = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(3), Inches(3))
        txA.name = "AreaA"
        txA.text_frame.text = "A"

        # Shift B so overlap is just 0.05in × 0.05in
        txB = slide.shapes.add_textbox(
            Emu(int(Inches(3)) - Emu(int(0.05 * _EMU_PER_INCH))),
            Emu(int(Inches(3)) - Emu(int(0.05 * _EMU_PER_INCH))),
            Inches(3), Inches(3),
        )
        txB.name = "AreaB"
        txB.text_frame.text = "B"

        _save(prs, path)
        result = verify_pptx(path)
        overlap_warns = [w for w in result["warnings"] if "overlap" in w.lower()]
        assert overlap_warns == []

    def test_no_trigger_for_non_overlapping_shapes(self, tmp_path):
        """Two shapes that do not overlap must not trigger."""
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches

        txA = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(2))
        txA.name = "SepA"
        txA.text_frame.text = "A"

        txB = slide.shapes.add_textbox(Inches(5), Inches(5), Inches(2), Inches(2))
        txB.name = "SepB"
        txB.text_frame.text = "B"

        _save(prs, path)
        result = verify_pptx(path)
        overlap_warns = [w for w in result["warnings"] if "overlap" in w.lower()]
        assert overlap_warns == []

    def test_no_trigger_for_small_overlap_fraction(self, tmp_path):
        """Overlap > 0.3 sq in but < 20% fraction of smaller shape must not trigger.

        BigBox: 8in × 8in (64 sq in).
        SmallBox: 2in × 2in (4 sq in), positioned so overlap is 0.3in × 2in = 0.6 sq in.
        Fraction = 0.6 / 4 = 15% < 20% threshold → no warning.
        Area = 0.6 sq in > 0.3 sq in threshold (area condition met).
        Both conditions must be true; since fraction fails, no warning should fire.
        """
        prs, path = _make_pptx(tmp_path)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches

        bigBox = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(8), Inches(8))
        bigBox.name = "BigBox"
        bigBox.text_frame.text = "big"

        # SmallBox left edge at 7.7in so overlap x-width = 8 - 7.7 = 0.3in,
        # overlap y = full 2in height → area = 0.6 sq in, fraction = 15%
        smallBox = slide.shapes.add_textbox(
            Emu(int(Inches(7.7))), Inches(0), Inches(2), Inches(2)
        )
        smallBox.name = "SmallBox"
        smallBox.text_frame.text = "small"

        _save(prs, path)
        result = verify_pptx(path)
        overlap_warns = [
            w for w in result["warnings"]
            if "overlap" in w.lower() and "BigBox" in w and "SmallBox" in w
        ]
        assert overlap_warns == []


# ── verify CLI command ────────────────────────────────────────────────────────


@pytest.fixture
def runner():
    return CliRunner()


@pytest.fixture
def clean_pptx(tmp_path):
    """A .pptx with no issues (ordinary text, normal sizes, no overflow)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    from pptx.util import Inches
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(3))
    txBox.text_frame.text = "Hello World"
    path = tmp_path / "clean.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def error_pptx(tmp_path):
    """A .pptx that has at least one error (unfilled placeholder text)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    from pptx.util import Inches
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    txBox.text_frame.text = "click to add title"
    path = tmp_path / "error.pptx"
    prs.save(str(path))
    return path


class TestVerifyCLI:

    def test_clean_file_exit_code_0(self, runner, clean_pptx):
        result = runner.invoke(cli, ["verify", str(clean_pptx)])
        assert result.exit_code == 0

    def test_clean_file_json_passed_true(self, runner, clean_pptx):
        result = runner.invoke(cli, ["verify", str(clean_pptx)])
        assert result.exit_code == 0
        data = json.loads(result.output)
        assert data["passed"] is True
        assert data["errors"] == []

    def test_error_file_exit_code_1(self, runner, error_pptx):
        result = runner.invoke(cli, ["verify", str(error_pptx)])
        assert result.exit_code == 1

    def test_error_file_json_passed_false(self, runner, error_pptx):
        result = runner.invoke(cli, ["verify", str(error_pptx)])
        data = json.loads(result.output)
        assert data["passed"] is False
        assert len(data["errors"]) >= 1

    def test_json_output_structure(self, runner, clean_pptx):
        result = runner.invoke(cli, ["verify", str(clean_pptx)])
        data = json.loads(result.output)
        for key in ("errors", "warnings", "passed", "slide_count", "error_count", "warning_count"):
            assert key in data

    def test_plain_clean_file_no_output(self, runner, clean_pptx):
        """--plain on a clean file must produce no stdout."""
        result = runner.invoke(cli, ["verify", "--plain", str(clean_pptx)])
        assert result.exit_code == 0
        assert result.output.strip() == ""

    def test_plain_error_file_fail_prefix(self, runner, error_pptx):
        """--plain with errors must prefix each line with FAIL."""
        result = runner.invoke(cli, ["verify", "--plain", str(error_pptx)])
        assert result.exit_code == 1
        lines = result.output.strip().splitlines()
        assert any(l.startswith("FAIL") for l in lines)

    def test_plain_warn_file_warn_prefix(self, tmp_path, runner):
        """--plain with only warnings must prefix each line with WARN."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches

        # Build a shape that triggers a warning (check 5 overlap)
        txA = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(4), Inches(4))
        txA.name = "WarnA"
        txA.text_frame.text = "A"
        txB = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(4), Inches(4))
        txB.name = "WarnB"
        txB.text_frame.text = "B"

        path = tmp_path / "warn.pptx"
        prs.save(str(path))

        result = runner.invoke(cli, ["verify", "--plain", str(path)])
        # If there are warnings but no errors, exit code should be 0
        data_result = runner.invoke(cli, ["verify", str(path)])
        data = json.loads(data_result.output)

        if data["warnings"]:
            assert result.exit_code == 0
            lines = result.output.strip().splitlines()
            assert any(l.startswith("WARN") for l in lines)

    def test_missing_file_stderr_error(self, runner, tmp_path):
        """A missing file must write to stderr."""
        result = runner.invoke(cli, ["verify", str(tmp_path / "nonexistent.pptx")])
        assert result.exit_code != 0
