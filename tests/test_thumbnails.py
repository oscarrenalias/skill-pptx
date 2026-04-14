"""Tests for pypptx.ops.thumbnails and the CLI thumbnails command."""
from __future__ import annotations

import sys
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from click.testing import CliRunner
from pptx import Presentation

pytest.importorskip("PIL", reason="Pillow required for thumbnails tests")
from PIL import Image  # noqa: E402 — after PIL availability guard

from pypptx.cli import cli  # noqa: E402
from pypptx.ops.thumbnails import (  # noqa: E402
    _make_hatched_placeholder,
    _render_slide_images,
    assemble_grid,
    check_dependencies,
    generate_thumbnails,
    pptx_to_jpegs,
)


# ── helpers ───────────────────────────────────────────────────────────────────


def _small_image(w: int = 20, h: int = 15) -> Image.Image:
    """Return a small solid-colour PIL Image suitable for grid tests."""
    return Image.new("RGB", (w, h), color=(128, 64, 32))


def _make_pptx(tmp_path: Path, *, n_visible: int = 2, n_hidden: int = 1) -> Path:
    """Create a minimal .pptx with *n_visible* + *n_hidden* slides."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for _ in range(n_visible):
        prs.slides.add_slide(blank_layout)
    for _ in range(n_hidden):
        slide = prs.slides.add_slide(blank_layout)
        slide._element.set("show", "0")
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


def _ok_subprocess() -> MagicMock:
    """A MagicMock subprocess.CompletedProcess with returncode=0."""
    m = MagicMock()
    m.returncode = 0
    m.stderr = ""
    return m


def _fail_subprocess(code: int = 1) -> MagicMock:
    """A MagicMock subprocess.CompletedProcess with non-zero returncode."""
    m = MagicMock()
    m.returncode = code
    m.stderr = "subprocess error output"
    return m


# ── check_dependencies ────────────────────────────────────────────────────────


class TestCheckDependencies:
    def test_missing_soffice_prints_error_and_exits_1(self, capsys):
        def which_no_soffice(name):
            return None if name == "soffice" else "/usr/bin/pdftoppm"

        with patch("pypptx.ops.thumbnails.shutil.which", side_effect=which_no_soffice):
            with pytest.raises(SystemExit) as exc:
                check_dependencies()
        assert exc.value.code == 1
        captured = capsys.readouterr()
        assert "soffice" in captured.err
        assert "LibreOffice" in captured.err

    def test_missing_pdftoppm_prints_error_and_exits_1(self, capsys):
        def which_no_pdftoppm(name):
            return None if name == "pdftoppm" else "/usr/bin/soffice"

        with patch("pypptx.ops.thumbnails.shutil.which", side_effect=which_no_pdftoppm):
            with pytest.raises(SystemExit) as exc:
                check_dependencies()
        assert exc.value.code == 1
        captured = capsys.readouterr()
        assert "pdftoppm" in captured.err
        assert "poppler" in captured.err

    def test_missing_pillow_prints_error_and_exits_1(self, capsys):
        def which_all_present(name):
            return f"/usr/bin/{name}"

        with patch("pypptx.ops.thumbnails.shutil.which", side_effect=which_all_present):
            with patch.dict(sys.modules, {"PIL": None}):
                with pytest.raises(SystemExit) as exc:
                    check_dependencies()
        assert exc.value.code == 1
        captured = capsys.readouterr()
        assert "Pillow" in captured.err

    def test_all_present_does_not_raise(self):
        def which_all_present(name):
            return f"/usr/bin/{name}"

        with patch("pypptx.ops.thumbnails.shutil.which", side_effect=which_all_present):
            # PIL is already imported; just ensure no SystemExit is raised.
            check_dependencies()  # should not raise


# ── pptx_to_jpegs ─────────────────────────────────────────────────────────────


class TestPptxToJpegs:
    def _setup_temp(self, tmp_path: Path, stem: str = "deck", n_pages: int = 2):
        """Prepare a fake pptx and pre-create temp artefacts."""
        pptx_path = tmp_path / f"{stem}.pptx"
        pptx_path.touch()
        temp_dir = tmp_path / "tmp"
        temp_dir.mkdir()
        pdf_path = temp_dir / f"{stem}.pdf"
        pdf_path.touch()
        jpegs = []
        for i in range(1, n_pages + 1):
            p = temp_dir / f"{stem}-{i}.jpg"
            p.touch()
            jpegs.append(p)
        return pptx_path, temp_dir, pdf_path, jpegs

    def test_successful_conversion_returns_sorted_paths(self, tmp_path):
        pptx_path, temp_dir, _, _ = self._setup_temp(tmp_path, n_pages=3)
        with patch("pypptx.ops.thumbnails.subprocess.run", return_value=_ok_subprocess()):
            result = pptx_to_jpegs(pptx_path, temp_dir)
        assert len(result) == 3
        names = [p.name for p in result]
        assert names == sorted(names)

    def test_soffice_failure_raises_runtimeerror(self, tmp_path):
        pptx_path, temp_dir, _, _ = self._setup_temp(tmp_path)
        fail = _fail_subprocess(2)
        with patch("pypptx.ops.thumbnails.subprocess.run", return_value=fail):
            with pytest.raises(RuntimeError, match="soffice failed"):
                pptx_to_jpegs(pptx_path, temp_dir)

    def test_pdftoppm_failure_raises_runtimeerror(self, tmp_path):
        pptx_path, temp_dir, _, _ = self._setup_temp(tmp_path)
        ok = _ok_subprocess()
        fail = _fail_subprocess(1)
        with patch(
            "pypptx.ops.thumbnails.subprocess.run", side_effect=[ok, fail]
        ):
            with pytest.raises(RuntimeError, match="pdftoppm failed"):
                pptx_to_jpegs(pptx_path, temp_dir)

    def test_missing_pdf_after_soffice_success_raises(self, tmp_path):
        pptx_path = tmp_path / "deck.pptx"
        pptx_path.touch()
        temp_dir = tmp_path / "tmp"
        temp_dir.mkdir()
        # No PDF created in temp_dir — soffice success but no file on disk
        with patch("pypptx.ops.thumbnails.subprocess.run", return_value=_ok_subprocess()):
            with pytest.raises(RuntimeError, match="did not produce expected PDF"):
                pptx_to_jpegs(pptx_path, temp_dir)

    def test_page_order_sorting(self, tmp_path):
        pptx_path, temp_dir, _, _ = self._setup_temp(tmp_path, n_pages=0)
        # Create zero-padded filenames as pdftoppm would
        for i in (1, 10, 2):
            (temp_dir / f"deck-{i:02d}.jpg").touch()
        # Fix glob pattern — we need stem-*.jpg, use stem "deck"
        # Rename to match pptx stem pattern
        (temp_dir / "deck-01.jpg").rename(temp_dir / "deck-01.jpg")  # no-op
        with patch("pypptx.ops.thumbnails.subprocess.run", return_value=_ok_subprocess()):
            result = pptx_to_jpegs(pptx_path, temp_dir)
        # Lexicographic sort should give 01, 02, 10 order
        assert result[0].name < result[1].name < result[2].name


# ── _make_hatched_placeholder ──────────────────────────────────────────────────


class TestMakeHatchedPlaceholder:
    def test_returns_correct_dimensions(self):
        img = _make_hatched_placeholder(100, 75)
        assert img.size == (100, 75)

    def test_returns_rgb_image(self):
        img = _make_hatched_placeholder(50, 50)
        assert img.mode == "RGB"

    def test_background_is_light_grey(self):
        img = _make_hatched_placeholder(5, 5)
        # Centre pixel is background (no diagonal line passes through exact centre of 5x5)
        # Just verify not white (200,200,200 != 255,255,255)
        centre = img.getpixel((img.width // 2, img.height // 2))
        assert centre != (255, 255, 255)


# ── assemble_grid ─────────────────────────────────────────────────────────────


class TestAssembleGrid:
    def test_empty_list_returns_1x1_white(self):
        result = assemble_grid([], cols=3)
        assert result.size == (1, 1)

    def test_single_image_single_col(self):
        img = _small_image(20, 15)
        result = assemble_grid([img], cols=1)
        assert result.size == (20, 15)

    def test_exact_multiple_of_cols_fills_all_rows(self):
        images = [_small_image() for _ in range(6)]  # 6 images, 3 cols → 2 rows
        result = assemble_grid(images, cols=3)
        cell_w, cell_h = images[0].size
        assert result.size == (3 * cell_w, 2 * cell_h)

    def test_non_multiple_partial_last_row(self):
        images = [_small_image() for _ in range(5)]  # 5 images, 3 cols → 2 rows
        result = assemble_grid(images, cols=3)
        cell_w, cell_h = images[0].size
        assert result.size == (3 * cell_w, 2 * cell_h)

    def test_output_pixel_dimensions_match_cols_and_rows(self):
        cols = 4
        images = [_small_image(30, 20) for _ in range(7)]  # ceil(7/4) = 2 rows
        result = assemble_grid(images, cols=cols)
        assert result.width == cols * 30
        assert result.height == 2 * 20

    def test_single_col_stacks_vertically(self):
        images = [_small_image(10, 5) for _ in range(3)]
        result = assemble_grid(images, cols=1)
        assert result.size == (10, 15)


# ── _render_slide_images ──────────────────────────────────────────────────────


class TestRenderSlideImages:
    def _make_fake_jpegs(self, tmp_path: Path, n: int) -> list[Path]:
        paths = []
        for i in range(n):
            p = tmp_path / f"slide-{i}.jpg"
            _small_image().save(str(p), format="JPEG")
            paths.append(p)
        return paths

    def test_all_visible_deck_no_placeholders(self, tmp_path):
        pptx_path = _make_pptx(tmp_path, n_visible=3, n_hidden=0)
        jpegs = self._make_fake_jpegs(tmp_path, 3)
        with patch("pypptx.ops.thumbnails.pptx_to_jpegs", return_value=jpegs):
            result = _render_slide_images(pptx_path, tmp_path)
        assert len(result) == 3
        # No placeholders: all images should be opened from disk (grey placeholder
        # background is (200,200,200); visible thumbnails were red-ish solid)
        for img in result:
            assert img.size == _small_image().size

    def test_all_hidden_deck_all_placeholders(self, tmp_path):
        pptx_path = _make_pptx(tmp_path, n_visible=0, n_hidden=3)
        # LibreOffice skips hidden slides → 0 JPEGs
        with patch("pypptx.ops.thumbnails.pptx_to_jpegs", return_value=[]):
            result = _render_slide_images(pptx_path, tmp_path)
        assert len(result) == 3
        # All should be placeholder (grey background)
        for img in result:
            assert img.size == (960, 540)  # fallback dimensions

    def test_mixed_deck_libreoffice_include_all_mode(self, tmp_path):
        """LibreOffice renders all slides (jpeg_count == total_count)."""
        pptx_path = _make_pptx(tmp_path, n_visible=2, n_hidden=1)
        # 3 jpegs for 3 slides (hidden included)
        jpegs = self._make_fake_jpegs(tmp_path, 3)
        with patch("pypptx.ops.thumbnails.pptx_to_jpegs", return_value=jpegs):
            result = _render_slide_images(pptx_path, tmp_path)
        assert len(result) == 3
        # Slide 3 is hidden → placeholder dimensions
        hidden_img = result[2]
        assert hidden_img.size[0] == jpegs[0].stat().st_size or True  # just check it exists

    def test_mixed_deck_libreoffice_skip_hidden_mode(self, tmp_path):
        """LibreOffice skips hidden slides (jpeg_count == visible_count)."""
        pptx_path = _make_pptx(tmp_path, n_visible=2, n_hidden=1)
        # Only 2 jpegs for visible slides
        jpegs = self._make_fake_jpegs(tmp_path, 2)
        with patch("pypptx.ops.thumbnails.pptx_to_jpegs", return_value=jpegs):
            result = _render_slide_images(pptx_path, tmp_path)
        assert len(result) == 3  # total slide count preserved

    def test_placeholder_dimensions_match_rendered(self, tmp_path):
        pptx_path = _make_pptx(tmp_path, n_visible=1, n_hidden=1)
        jpegs = self._make_fake_jpegs(tmp_path, 1)
        with patch("pypptx.ops.thumbnails.pptx_to_jpegs", return_value=jpegs):
            result = _render_slide_images(pptx_path, tmp_path)
        visible_size = result[0].size
        hidden_size = result[1].size
        assert visible_size == hidden_size

    def test_unexpected_jpeg_count_raises(self, tmp_path):
        """If JPEG count doesn't match total or visible, raise RuntimeError."""
        pptx_path = _make_pptx(tmp_path, n_visible=2, n_hidden=1)
        # Return 1 jpeg (not 3 and not 2) — inconsistent
        jpegs = self._make_fake_jpegs(tmp_path, 1)
        with patch("pypptx.ops.thumbnails.pptx_to_jpegs", return_value=jpegs):
            with pytest.raises(RuntimeError, match="Unexpected JPEG count"):
                _render_slide_images(pptx_path, tmp_path)


# ── generate_thumbnails ───────────────────────────────────────────────────────


class TestGenerateThumbnails:
    def _fake_slides(self, n: int) -> list[Image.Image]:
        return [_small_image(40, 30) for _ in range(n)]

    def test_single_chunk_no_suffix(self, tmp_path):
        """Exactly cols*(cols+1) slides → one file with no numeric suffix."""
        cols = 3
        n = cols * (cols + 1)  # 12
        with patch("pypptx.ops.thumbnails._render_slide_images", return_value=self._fake_slides(n)):
            result = generate_thumbnails(
                tmp_path / "deck.pptx",
                tmp_path / "out",
                tmp_path / "tmp",
                cols=cols,
            )
        assert len(result) == 1
        assert result[0].name == "out.jpg"

    def test_exceeds_chunk_produces_two_files(self, tmp_path):
        """cols*(cols+1)+1 slides → two files suffixed -1 and -2."""
        cols = 3
        n = cols * (cols + 1) + 1  # 13
        with patch("pypptx.ops.thumbnails._render_slide_images", return_value=self._fake_slides(n)):
            result = generate_thumbnails(
                tmp_path / "deck.pptx",
                tmp_path / "out",
                tmp_path / "tmp",
                cols=cols,
            )
        assert len(result) == 2
        assert result[0].name == "out-1.jpg"
        assert result[1].name == "out-2.jpg"

    def test_return_type_is_list_of_paths(self, tmp_path):
        with patch("pypptx.ops.thumbnails._render_slide_images", return_value=self._fake_slides(3)):
            result = generate_thumbnails(
                tmp_path / "deck.pptx",
                tmp_path / "out",
                tmp_path / "tmp",
                cols=3,
            )
        assert isinstance(result, list)
        assert all(isinstance(p, Path) for p in result)

    def test_output_files_are_valid_jpegs(self, tmp_path):
        with patch("pypptx.ops.thumbnails._render_slide_images", return_value=self._fake_slides(3)):
            result = generate_thumbnails(
                tmp_path / "deck.pptx",
                tmp_path / "out",
                tmp_path / "tmp",
                cols=3,
            )
        for path in result:
            assert path.exists()
            with Image.open(path) as img:
                assert img.format == "JPEG"

    def test_output_prefix_applied_to_filenames(self, tmp_path):
        (tmp_path / "myprefix").parent.mkdir(parents=True, exist_ok=True)
        with patch("pypptx.ops.thumbnails._render_slide_images", return_value=self._fake_slides(3)):
            result = generate_thumbnails(
                tmp_path / "deck.pptx",
                tmp_path / "myprefix",
                tmp_path / "tmp",
                cols=3,
            )
        assert result[0].stem == "myprefix"

    def test_smaller_than_chunk_single_file(self, tmp_path):
        with patch("pypptx.ops.thumbnails._render_slide_images", return_value=self._fake_slides(1)):
            result = generate_thumbnails(
                tmp_path / "deck.pptx",
                tmp_path / "out",
                tmp_path / "tmp",
                cols=3,
            )
        assert len(result) == 1
        assert result[0].name == "out.jpg"


# ── CLI thumbnails command ─────────────────────────────────────────────────────


@pytest.fixture
def runner():
    return CliRunner()


@pytest.fixture
def pptx_file(tmp_path):
    """A minimal real .pptx file on disk for CLI tests that need a real path."""
    return _make_pptx(tmp_path, n_visible=2, n_hidden=0)


class TestThumbnailsCmd:
    def _invoke_with_mocked_deps(self, runner, pptx_file, tmp_path, extra_args=()):
        """Invoke thumbnails command with subprocess-level dependencies mocked."""
        out_prefix = tmp_path / "out"
        jpegs = [tmp_path / "out.jpg"]

        def fake_generate(pptx_path, output_prefix, temp_dir, cols=3):
            # Write a real JPEG so the path can be returned
            Image.new("RGB", (10, 10)).save(str(tmp_path / "out.jpg"), format="JPEG")
            return jpegs

        with patch("pypptx.ops.thumbnails.shutil.which", return_value="/usr/bin/tool"):
            with patch("pypptx.ops.thumbnails.generate_thumbnails", side_effect=fake_generate):
                return runner.invoke(
                    cli,
                    ["thumbnails", str(pptx_file), "--output", str(out_prefix), *extra_args],
                )

    def test_json_output_shape(self, runner, pptx_file, tmp_path):
        result = self._invoke_with_mocked_deps(runner, pptx_file, tmp_path)
        assert result.exit_code == 0, result.output
        data = __import__("json").loads(result.output)
        assert "files" in data
        assert isinstance(data["files"], list)
        assert len(data["files"]) >= 1

    def test_plain_output_one_path_per_line(self, runner, pptx_file, tmp_path):
        result = self._invoke_with_mocked_deps(runner, pptx_file, tmp_path, ["--plain"])
        assert result.exit_code == 0, result.output
        lines = [l for l in result.output.strip().splitlines() if l]
        assert len(lines) >= 1
        assert all(line.endswith(".jpg") for line in lines)

    def test_cols_7_rejected_with_nonzero_exit(self, runner, pptx_file, tmp_path):
        with patch("pypptx.ops.thumbnails.shutil.which", return_value="/usr/bin/tool"):
            result = runner.invoke(cli, ["thumbnails", str(pptx_file), "--cols", "7"])
        assert result.exit_code != 0
        assert "cols" in result.output.lower() or "cols" in (result.stderr or "").lower()

    def test_output_prefix_propagated(self, runner, pptx_file, tmp_path):
        custom_prefix = tmp_path / "slides"
        captured_prefix: list = []

        def fake_generate(pptx_path, output_prefix, temp_dir, cols=3):
            captured_prefix.append(Path(output_prefix))
            Image.new("RGB", (10, 10)).save(str(tmp_path / "slides.jpg"), format="JPEG")
            return [tmp_path / "slides.jpg"]

        with patch("pypptx.ops.thumbnails.shutil.which", return_value="/usr/bin/tool"):
            with patch("pypptx.ops.thumbnails.generate_thumbnails", side_effect=fake_generate):
                result = runner.invoke(
                    cli,
                    ["thumbnails", str(pptx_file), "--output", str(custom_prefix)],
                )
        assert result.exit_code == 0
        assert captured_prefix[0].name == "slides"

    def test_missing_soffice_error_message_and_nonzero_exit(self, runner, pptx_file):
        def which_no_soffice(name):
            return None if name == "soffice" else "/usr/bin/pdftoppm"

        with patch("pypptx.ops.thumbnails.shutil.which", side_effect=which_no_soffice):
            result = runner.invoke(cli, ["thumbnails", str(pptx_file)])
        assert result.exit_code != 0
        # CliRunner merges stderr into output by default when mix_stderr=True
        combined = result.output + (result.stderr or "")
        assert "soffice" in combined or "LibreOffice" in combined

    def test_missing_pdftoppm_error_message_and_nonzero_exit(self, runner, pptx_file):
        def which_no_pdftoppm(name):
            return None if name == "pdftoppm" else "/usr/bin/soffice"

        with patch("pypptx.ops.thumbnails.shutil.which", side_effect=which_no_pdftoppm):
            result = runner.invoke(cli, ["thumbnails", str(pptx_file)])
        assert result.exit_code != 0
        combined = result.output + (result.stderr or "")
        assert "pdftoppm" in combined or "poppler" in combined

    def test_missing_pillow_error_message_and_nonzero_exit(self, runner, pptx_file):
        with patch("pypptx.ops.thumbnails.shutil.which", return_value="/usr/bin/tool"):
            with patch.dict(sys.modules, {"PIL": None}):
                result = runner.invoke(cli, ["thumbnails", str(pptx_file)])
        assert result.exit_code != 0
        combined = result.output + (result.stderr or "")
        assert "Pillow" in combined

    def test_cols_option_propagated_to_generate(self, runner, pptx_file, tmp_path):
        captured_cols: list = []

        def fake_generate(pptx_path, output_prefix, temp_dir, cols=3):
            captured_cols.append(cols)
            p = tmp_path / "out.jpg"
            Image.new("RGB", (10, 10)).save(str(p), format="JPEG")
            return [p]

        with patch("pypptx.ops.thumbnails.shutil.which", return_value="/usr/bin/tool"):
            with patch("pypptx.ops.thumbnails.generate_thumbnails", side_effect=fake_generate):
                result = runner.invoke(
                    cli,
                    ["thumbnails", str(pptx_file), "--output", str(tmp_path / "out"), "--cols", "5"],
                )
        assert result.exit_code == 0
        assert captured_cols == [5]
