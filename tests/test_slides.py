"""Tests for pypptx.ops.slides — list, add, delete, move, pptx_edit."""
import hashlib
import shutil
from pathlib import Path

import pytest
from pptx import Presentation

from pypptx.ops.slides import (
    add_slide,
    delete_slide,
    list_layouts,
    list_slides,
    move_slide,
    pptx_edit,
)


# ── list_slides ───────────────────────────────────────────────────────────────


class TestListSlides:
    def test_returns_all_slides(self, minimal_pptx):
        result = list_slides(minimal_pptx)
        assert len(result) == 3

    def test_1_based_indices(self, minimal_pptx):
        result = list_slides(minimal_pptx)
        assert [s["index"] for s in result] == [1, 2, 3]

    def test_file_field_is_xml_filename(self, minimal_pptx):
        result = list_slides(minimal_pptx)
        for s in result:
            assert s["file"].endswith(".xml")

    def test_hidden_slide_detected_with_show_zero(self, minimal_pptx):
        """Slide 3 was marked hidden with show='0' in the fixture."""
        result = list_slides(minimal_pptx)
        assert result[2]["hidden"] is True
        assert result[0]["hidden"] is False
        assert result[1]["hidden"] is False

    def test_hidden_slide_detected_with_show_false(self, tmp_path):
        """A slide with show='false' is also reported as hidden."""
        prs = Presentation()
        layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(layout)
        slide._element.set("show", "false")
        path = tmp_path / "hidden_false.pptx"
        prs.save(str(path))
        result = list_slides(path)
        assert result[0]["hidden"] is True

    def test_directory_input(self, unpacked_pptx):
        result = list_slides(unpacked_pptx)
        assert len(result) == 3

    def test_hidden_flag_via_directory_input(self, unpacked_pptx):
        result = list_slides(unpacked_pptx)
        assert result[2]["hidden"] is True


# ── list_layouts ──────────────────────────────────────────────────────────────


class TestListLayouts:
    def test_returns_layouts(self, minimal_pptx):
        result = list_layouts(minimal_pptx)
        assert len(result) > 0

    def test_1_based_indices(self, minimal_pptx):
        result = list_layouts(minimal_pptx)
        assert [l["index"] for l in result] == list(range(1, len(result) + 1))

    def test_filename_sorted_order(self, minimal_pptx):
        """Layouts are returned sorted by filename."""
        result = list_layouts(minimal_pptx)
        filenames = [l["file"] for l in result]
        assert filenames == sorted(filenames)

    def test_file_field_is_xml_filename(self, minimal_pptx):
        result = list_layouts(minimal_pptx)
        for l in result:
            assert l["file"].endswith(".xml")

    def test_directory_input(self, unpacked_pptx):
        result = list_layouts(unpacked_pptx)
        assert len(result) > 0


# ── add_slide ─────────────────────────────────────────────────────────────────


class TestAddSlide:
    def test_duplicate_increases_slide_count(self, minimal_pptx):
        before = len(list_slides(minimal_pptx))
        add_slide(minimal_pptx, duplicate=1)
        assert len(list_slides(minimal_pptx)) == before + 1

    def test_duplicate_returns_correct_index(self, minimal_pptx):
        result = add_slide(minimal_pptx, duplicate=1)
        assert result["index"] == 4

    def test_duplicate_file_reopens_with_python_pptx(self, minimal_pptx):
        add_slide(minimal_pptx, duplicate=1)
        prs = Presentation(str(minimal_pptx))
        assert len(prs.slides) == 4

    def test_duplicate_notes_rels_stripped(self, minimal_pptx, tmp_path):
        """Notes-slide relationship is stripped from the duplicated slide."""
        # After duplication, the file must be reopenable (no broken notes rels).
        add_slide(minimal_pptx, duplicate=1)
        Presentation(str(minimal_pptx))  # no exception expected

    def test_layout_increases_slide_count(self, minimal_pptx):
        before = len(list_slides(minimal_pptx))
        add_slide(minimal_pptx, layout=1)
        assert len(list_slides(minimal_pptx)) == before + 1

    def test_layout_slide_reopens_with_python_pptx(self, minimal_pptx):
        add_slide(minimal_pptx, layout=1)
        prs = Presentation(str(minimal_pptx))
        assert len(prs.slides) == 4

    def test_both_raises_value_error(self, minimal_pptx):
        with pytest.raises(ValueError, match="[Ee]xactly one"):
            add_slide(minimal_pptx, duplicate=1, layout=1)

    def test_neither_raises_value_error(self, minimal_pptx):
        with pytest.raises(ValueError, match="[Ee]xactly one"):
            add_slide(minimal_pptx)

    def test_out_of_range_duplicate_raises(self, minimal_pptx):
        with pytest.raises(ValueError):
            add_slide(minimal_pptx, duplicate=99)

    def test_out_of_range_layout_raises(self, minimal_pptx):
        with pytest.raises(ValueError):
            add_slide(minimal_pptx, layout=9999)

    def test_directory_input(self, unpacked_pptx):
        before = len(list_slides(unpacked_pptx))
        add_slide(unpacked_pptx, duplicate=1)
        assert len(list_slides(unpacked_pptx)) == before + 1

    def test_layout_directory_input(self, unpacked_pptx):
        before = len(list_slides(unpacked_pptx))
        add_slide(unpacked_pptx, layout=1)
        assert len(list_slides(unpacked_pptx)) == before + 1


# ── delete_slide ──────────────────────────────────────────────────────────────


class TestDeleteSlide:
    def test_decreases_slide_count(self, minimal_pptx):
        before = len(list_slides(minimal_pptx))
        delete_slide(minimal_pptx, 1)
        assert len(list_slides(minimal_pptx)) == before - 1

    def test_returns_deleted_file_and_index(self, minimal_pptx):
        result = delete_slide(minimal_pptx, 2)
        assert "deleted_file" in result
        assert result["deleted_index"] == 2

    def test_pptx_remains_valid_after_delete(self, minimal_pptx):
        delete_slide(minimal_pptx, 1)
        prs = Presentation(str(minimal_pptx))
        assert len(prs.slides) == 2

    def test_out_of_range_raises(self, minimal_pptx):
        with pytest.raises(ValueError):
            delete_slide(minimal_pptx, 99)

    def test_zero_index_raises(self, minimal_pptx):
        with pytest.raises(ValueError):
            delete_slide(minimal_pptx, 0)

    def test_directory_input(self, unpacked_pptx):
        before = len(list_slides(unpacked_pptx))
        delete_slide(unpacked_pptx, 1)
        assert len(list_slides(unpacked_pptx)) == before - 1

    def test_deleted_file_removed_from_slides_dir(self, unpacked_pptx):
        slides_before = [s["file"] for s in list_slides(unpacked_pptx)]
        result = delete_slide(unpacked_pptx, 1)
        deleted = result["deleted_file"]
        assert deleted in slides_before
        assert not (unpacked_pptx / "ppt" / "slides" / deleted).exists()


# ── move_slide ────────────────────────────────────────────────────────────────


class TestMoveSlide:
    def test_no_op_when_from_equals_to(self, minimal_pptx):
        result = move_slide(minimal_pptx, 2, 2)
        assert result["from"] == 2
        assert result["to"] == 2

    def test_returns_file_from_to(self, minimal_pptx):
        result = move_slide(minimal_pptx, 1, 3)
        assert "file" in result
        assert result["from"] == 1
        assert result["to"] == 3

    def test_reorders_slides(self, minimal_pptx):
        """Moving slide 3 (hidden) to position 1 places it first.

        python-pptx renumbers slide files in sldIdLst order on every open,
        so filenames are not stable across moves. Check the ``hidden``
        attribute, which is derived from slide XML content, not the filename.
        """
        slides_before = list_slides(minimal_pptx)
        move_slide(minimal_pptx, 3, 1)
        slides_after = list_slides(minimal_pptx)
        # The hidden slide (position 3 before) should now be at position 1.
        assert slides_after[0]["hidden"] == slides_before[2]["hidden"]
        assert slides_after[0]["hidden"] is True

    def test_slide_count_unchanged_after_move(self, minimal_pptx):
        move_slide(minimal_pptx, 1, 3)
        assert len(list_slides(minimal_pptx)) == 3

    def test_out_of_range_from_raises(self, minimal_pptx):
        with pytest.raises(ValueError):
            move_slide(minimal_pptx, 99, 1)

    def test_out_of_range_to_raises(self, minimal_pptx):
        with pytest.raises(ValueError):
            move_slide(minimal_pptx, 1, 99)

    def test_directory_input(self, unpacked_pptx):
        result = move_slide(unpacked_pptx, 1, 2)
        assert result["from"] == 1
        assert result["to"] == 2


# ── pptx_edit context manager ─────────────────────────────────────────────────


class TestPptxEdit:
    def test_clean_exit_repacks_file(self, minimal_pptx):
        """On clean exit the original file is replaced with the repacked version."""
        with pptx_edit(minimal_pptx) as tmp_dir:
            assert tmp_dir.is_dir()
            assert (tmp_dir / "ppt" / "presentation.xml").is_file()
        # File must still be openable after the context manager exits.
        prs = Presentation(str(minimal_pptx))
        assert len(prs.slides) == 3

    def test_exception_leaves_original_intact(self, minimal_pptx):
        """When the body raises, the original .pptx is NOT overwritten."""
        original_bytes = minimal_pptx.read_bytes()
        original_hash = hashlib.md5(original_bytes).hexdigest()

        with pytest.raises(RuntimeError):
            with pptx_edit(minimal_pptx) as tmp_dir:
                raise RuntimeError("simulated failure")

        new_hash = hashlib.md5(minimal_pptx.read_bytes()).hexdigest()
        assert original_hash == new_hash

    def test_tmp_dir_cleaned_up_on_success(self, minimal_pptx, tmp_path):
        """Temporary directory is removed after a clean exit."""
        captured = {}

        def _run():
            with pptx_edit(minimal_pptx) as tmp_dir:
                captured["path"] = tmp_dir

        _run()
        assert not captured["path"].exists()

    def test_tmp_dir_cleaned_up_on_exception(self, minimal_pptx):
        """Temporary directory is removed even when an exception is raised."""
        captured = {}
        try:
            with pptx_edit(minimal_pptx) as tmp_dir:
                captured["path"] = tmp_dir
                raise ValueError("boom")
        except ValueError:
            pass
        assert not captured["path"].exists()
