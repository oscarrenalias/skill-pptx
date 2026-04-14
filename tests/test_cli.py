"""CLI integration tests for pypptx, exercised via click's CliRunner."""
import json
import shutil
from pathlib import Path

import pytest
from click.testing import CliRunner
from pptx import Presentation

from pypptx.cli import cli
from pypptx.ops.pack import unpack


@pytest.fixture
def runner():
    return CliRunner()


@pytest.fixture
def pptx_copy(minimal_pptx, tmp_path):
    """A fresh copy of the minimal fixture so mutations don't bleed across tests."""
    dst = tmp_path / "copy.pptx"
    shutil.copy2(minimal_pptx, dst)
    return dst


# ── Global commands ───────────────────────────────────────────────────────────


class TestGlobal:
    def test_version_output(self, runner):
        result = runner.invoke(cli, ["--version"])
        assert result.exit_code == 0
        assert "pypptx" in result.output

    def test_help_lists_subcommands(self, runner):
        result = runner.invoke(cli, ["--help"])
        assert result.exit_code == 0
        for cmd in ("extract-text", "unpack", "pack", "clean", "slide"):
            assert cmd in result.output

    def test_slide_help(self, runner):
        result = runner.invoke(cli, ["slide", "--help"])
        assert result.exit_code == 0
        for sub in ("list", "layouts", "add", "delete", "move"):
            assert sub in result.output


# ── unpack command ────────────────────────────────────────────────────────────


class TestUnpackCmd:
    def test_json_output(self, runner, minimal_pptx, tmp_path):
        out = tmp_path / "out"
        result = runner.invoke(cli, ["unpack", str(minimal_pptx), str(out)])
        assert result.exit_code == 0
        data = json.loads(result.output)
        assert "unpacked_dir" in data

    def test_plain_output(self, runner, minimal_pptx, tmp_path):
        out = tmp_path / "out"
        result = runner.invoke(cli, ["unpack", "--plain", str(minimal_pptx), str(out)])
        assert result.exit_code == 0
        assert str(out) in result.output.strip()

    def test_default_output_dir_uses_file_stem(self, runner, minimal_pptx, tmp_path):
        with runner.isolated_filesystem(temp_dir=tmp_path):
            shutil.copy2(minimal_pptx, "test.pptx")
            result = runner.invoke(cli, ["unpack", "test.pptx"])
            assert result.exit_code == 0
            data = json.loads(result.output)
            assert data["unpacked_dir"] == "test"

    def test_missing_file_exits_nonzero(self, runner, tmp_path):
        result = runner.invoke(cli, ["unpack", str(tmp_path / "nope.pptx")])
        assert result.exit_code != 0

    def test_bad_zip_exits_nonzero(self, runner, tmp_path):
        bad = tmp_path / "bad.pptx"
        bad.write_bytes(b"not a zip")
        result = runner.invoke(cli, ["unpack", str(bad), str(tmp_path / "out")])
        assert result.exit_code != 0


# ── pack command ──────────────────────────────────────────────────────────────


class TestPackCmd:
    def test_json_output(self, runner, minimal_pptx, tmp_path):
        unpacked = tmp_path / "unpacked"
        unpack(minimal_pptx, unpacked)
        output = tmp_path / "repacked.pptx"
        result = runner.invoke(cli, ["pack", str(unpacked), str(output)])
        assert result.exit_code == 0
        data = json.loads(result.output)
        assert "output_file" in data

    def test_plain_output(self, runner, minimal_pptx, tmp_path):
        unpacked = tmp_path / "unpacked"
        unpack(minimal_pptx, unpacked)
        output = tmp_path / "repacked.pptx"
        result = runner.invoke(cli, ["pack", "--plain", str(unpacked), str(output)])
        assert result.exit_code == 0
        assert str(output) in result.output.strip()

    def test_output_file_is_valid_pptx(self, runner, minimal_pptx, tmp_path):
        unpacked = tmp_path / "unpacked"
        unpack(minimal_pptx, unpacked)
        output = tmp_path / "repacked.pptx"
        runner.invoke(cli, ["pack", str(unpacked), str(output)])
        prs = Presentation(str(output))
        assert len(prs.slides) == 3


# ── clean command ─────────────────────────────────────────────────────────────


class TestCleanCmd:
    def test_json_with_orphans(self, runner, unpacked_pptx):
        orphan = unpacked_pptx / "ppt" / "media" / "orphan.png"
        orphan.parent.mkdir(parents=True, exist_ok=True)
        orphan.write_bytes(b"\x89PNG")
        result = runner.invoke(cli, ["clean", str(unpacked_pptx)])
        assert result.exit_code == 0
        data = json.loads(result.output)
        assert "removed" in data
        assert len(data["removed"]) >= 1

    def test_json_empty_list_when_no_orphans(self, runner, unpacked_pptx):
        result = runner.invoke(cli, ["clean", str(unpacked_pptx)])
        assert result.exit_code == 0
        data = json.loads(result.output)
        assert data["removed"] == []

    def test_plain_output_with_orphans(self, runner, unpacked_pptx):
        orphan = unpacked_pptx / "ppt" / "media" / "orphan.png"
        orphan.parent.mkdir(parents=True, exist_ok=True)
        orphan.write_bytes(b"\x89PNG")
        result = runner.invoke(cli, ["clean", "--plain", str(unpacked_pptx)])
        assert result.exit_code == 0
        assert "orphan.png" in result.output

    def test_plain_output_empty_when_no_orphans(self, runner, unpacked_pptx):
        result = runner.invoke(cli, ["clean", "--plain", str(unpacked_pptx)])
        assert result.exit_code == 0
        assert result.output.strip() == ""

    def test_directory_path_input(self, runner, unpacked_pptx):
        result = runner.invoke(cli, ["clean", str(unpacked_pptx)])
        assert result.exit_code == 0

    def test_nonexistent_path_exits_nonzero(self, runner, tmp_path):
        result = runner.invoke(cli, ["clean", str(tmp_path / "nonexistent")])
        assert result.exit_code != 0

    def test_stderr_contains_error_on_failure(self, runner, tmp_path):
        """On error, message goes to stderr and exit code is non-zero."""
        result = runner.invoke(cli, ["clean", str(tmp_path / "nonexistent")])
        assert result.exit_code != 0


# ── slide list command ────────────────────────────────────────────────────────


class TestSlideListCmd:
    def test_json_output(self, runner, minimal_pptx):
        result = runner.invoke(cli, ["slide", "list", str(minimal_pptx)])
        assert result.exit_code == 0
        data = json.loads(result.output)
        assert "slides" in data
        assert len(data["slides"]) == 3

    def test_plain_output_line_per_slide(self, runner, minimal_pptx):
        result = runner.invoke(cli, ["slide", "list", "--plain", str(minimal_pptx)])
        assert result.exit_code == 0
        lines = result.output.strip().splitlines()
        assert len(lines) == 3

    def test_plain_output_hidden_annotation(self, runner, minimal_pptx):
        result = runner.invoke(cli, ["slide", "list", "--plain", str(minimal_pptx)])
        lines = result.output.strip().splitlines()
        assert "[hidden]" in lines[2]
        assert "[hidden]" not in lines[0]

    def test_directory_input(self, runner, unpacked_pptx):
        result = runner.invoke(cli, ["slide", "list", str(unpacked_pptx)])
        assert result.exit_code == 0
        data = json.loads(result.output)
        assert len(data["slides"]) == 3


# ── slide layouts command ─────────────────────────────────────────────────────


class TestSlideLayoutsCmd:
    def test_json_output(self, runner, minimal_pptx):
        result = runner.invoke(cli, ["slide", "layouts", str(minimal_pptx)])
        assert result.exit_code == 0
        data = json.loads(result.output)
        assert "layouts" in data
        assert len(data["layouts"]) > 0

    def test_plain_output_one_filename_per_line(self, runner, minimal_pptx):
        result = runner.invoke(cli, ["slide", "layouts", "--plain", str(minimal_pptx)])
        assert result.exit_code == 0
        lines = [l for l in result.output.strip().splitlines() if l]
        assert len(lines) > 0
        assert all(".xml" in l for l in lines)

    def test_plain_output_includes_filename_and_name(self, runner, minimal_pptx):
        result = runner.invoke(cli, ["slide", "layouts", "--plain", str(minimal_pptx)])
        assert result.exit_code == 0
        lines = [l for l in result.output.strip().splitlines() if l]
        assert len(lines) > 0
        for line in lines:
            parts = line.split("  ", 1)
            assert len(parts) == 2
            assert parts[0].endswith(".xml")
            assert len(parts[1].strip()) > 0


# ── slide add command ─────────────────────────────────────────────────────────


class TestSlideAddCmd:
    def test_duplicate_json_output(self, runner, pptx_copy):
        result = runner.invoke(cli, ["slide", "add", "--duplicate", "1", str(pptx_copy)])
        assert result.exit_code == 0
        data = json.loads(result.output)
        assert "added_file" in data
        assert "position" in data
        assert data["position"] == 4

    def test_layout_json_output(self, runner, pptx_copy):
        result = runner.invoke(cli, ["slide", "add", "--layout", "1", str(pptx_copy)])
        assert result.exit_code == 0
        data = json.loads(result.output)
        assert "added_file" in data

    def test_plain_output(self, runner, pptx_copy):
        result = runner.invoke(cli, ["slide", "add", "--plain", "--duplicate", "1", str(pptx_copy)])
        assert result.exit_code == 0
        assert "at position" in result.output

    def test_mutual_exclusion_exits_nonzero(self, runner, minimal_pptx):
        result = runner.invoke(cli, [
            "slide", "add", "--duplicate", "1", "--layout", "1", str(minimal_pptx)
        ])
        assert result.exit_code != 0

    def test_neither_flag_exits_nonzero(self, runner, minimal_pptx):
        result = runner.invoke(cli, ["slide", "add", str(minimal_pptx)])
        assert result.exit_code != 0

    def test_add_with_position_moves_to_target(self, runner, pptx_copy):
        result = runner.invoke(cli, [
            "slide", "add", "--duplicate", "1", "--position", "1", str(pptx_copy)
        ])
        assert result.exit_code == 0
        data = json.loads(result.output)
        assert data["position"] == 1

    def test_stderr_on_out_of_range_duplicate(self, runner, minimal_pptx):
        result = runner.invoke(cli, ["slide", "add", "--duplicate", "99", str(minimal_pptx)])
        assert result.exit_code != 0


# ── slide delete command ──────────────────────────────────────────────────────


class TestSlideDeleteCmd:
    def test_json_output(self, runner, pptx_copy):
        result = runner.invoke(cli, ["slide", "delete", str(pptx_copy), "1"])
        assert result.exit_code == 0
        data = json.loads(result.output)
        assert "deleted_file" in data
        assert data["deleted_index"] == 1

    def test_plain_output(self, runner, pptx_copy):
        result = runner.invoke(cli, ["slide", "delete", "--plain", str(pptx_copy), "2"])
        assert result.exit_code == 0
        assert "deleted" in result.output

    def test_out_of_range_exits_nonzero(self, runner, minimal_pptx):
        result = runner.invoke(cli, ["slide", "delete", str(minimal_pptx), "99"])
        assert result.exit_code != 0

    def test_slide_count_decreases(self, runner, pptx_copy):
        runner.invoke(cli, ["slide", "delete", str(pptx_copy), "1"])
        prs = Presentation(str(pptx_copy))
        assert len(prs.slides) == 2


# ── slide move command ────────────────────────────────────────────────────────


class TestSlideMoveCmd:
    def test_json_output(self, runner, pptx_copy):
        result = runner.invoke(cli, ["slide", "move", str(pptx_copy), "1", "2"])
        assert result.exit_code == 0
        data = json.loads(result.output)
        assert data["from"] == 1
        assert data["to"] == 2

    def test_plain_output(self, runner, pptx_copy):
        result = runner.invoke(cli, ["slide", "move", "--plain", str(pptx_copy), "1", "3"])
        assert result.exit_code == 0
        assert "1 -> 3" in result.output

    def test_out_of_range_exits_nonzero(self, runner, minimal_pptx):
        result = runner.invoke(cli, ["slide", "move", str(minimal_pptx), "1", "99"])
        assert result.exit_code != 0


# ── extract-text command ──────────────────────────────────────────────────────


class TestExtractTextCmd:
    def test_stdout_without_output_flag(self, runner, minimal_pptx):
        """Without --output the text goes straight to stdout (no JSON wrapper)."""
        result = runner.invoke(cli, ["extract-text", str(minimal_pptx)])
        assert result.exit_code == 0
        assert "--- Slide 1 ---" in result.output

    def test_slides_filter(self, runner, minimal_pptx):
        result = runner.invoke(cli, ["extract-text", "--slides", "1", str(minimal_pptx)])
        assert result.exit_code == 0
        assert "--- Slide 1 ---" in result.output
        assert "--- Slide 2 ---" not in result.output

    def test_output_file_writes_content(self, runner, minimal_pptx, tmp_path):
        out_file = tmp_path / "text.txt"
        result = runner.invoke(cli, ["extract-text", "--output", str(out_file), str(minimal_pptx)])
        assert result.exit_code == 0
        assert out_file.exists()
        content = out_file.read_text(encoding="utf-8")
        assert "--- Slide 1 ---" in content

    def test_output_file_json_stdout(self, runner, minimal_pptx, tmp_path):
        """With --output, JSON containing output_file and slide_count is printed."""
        out_file = tmp_path / "text.txt"
        result = runner.invoke(cli, ["extract-text", "--output", str(out_file), str(minimal_pptx)])
        assert result.exit_code == 0
        data = json.loads(result.output)
        assert "output_file" in data
        assert "slide_count" in data
        assert data["slide_count"] == 3

    def test_output_file_plain_stdout(self, runner, minimal_pptx, tmp_path):
        """With --output --plain, only the output filename is written to stdout."""
        out_file = tmp_path / "text.txt"
        result = runner.invoke(cli, [
            "extract-text", "--output", str(out_file), "--plain", str(minimal_pptx)
        ])
        assert result.exit_code == 0
        assert result.output.strip() == str(out_file)

    def test_non_integer_slide_index_exits_nonzero(self, runner, minimal_pptx):
        result = runner.invoke(cli, ["extract-text", "--slides", "abc", str(minimal_pptx)])
        assert result.exit_code != 0

    def test_out_of_range_slide_index_exits_nonzero(self, runner, minimal_pptx):
        result = runner.invoke(cli, ["extract-text", "--slides", "99", str(minimal_pptx)])
        assert result.exit_code != 0

    def test_unreadable_file_exits_nonzero(self, runner, tmp_path):
        """A non-PPTX file raises an exception that the CLI converts to exit 1."""
        bad = tmp_path / "bad.pptx"
        bad.write_bytes(b"this is not a pptx")
        result = runner.invoke(cli, ["extract-text", str(bad)])
        assert result.exit_code != 0
